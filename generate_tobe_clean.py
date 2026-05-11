# -*- coding: utf-8 -*-
"""
TeleResolve - To-Be Diagram (No X circles, white background)
Same KPMG swim-lane format. X circles removed. All backgrounds white.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

NAVY  = RGBColor(0x1F, 0x39, 0x64)
BLUE  = RGBColor(0x41, 0x72, 0xC4)
CYAN  = RGBColor(0x00, 0xB0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY  = RGBColor(0xBF, 0xBF, 0xBF)
DGREY = RGBColor(0x35, 0x35, 0x35)
RED   = RGBColor(0xC0, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
sl = prs.slides.add_slide(prs.slide_layouts[6])

def _sp(mso, x, y, w, h, fill, lc=GREY, lw=Pt(0.75)):
    s = sl.shapes.add_shape(mso, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc is None: s.line.fill.background()
    else: s.line.color.rgb = lc; s.line.width = lw
    s.text_frame.word_wrap = True
    return s

def RCT(x,y,w,h, fill=WHITE, lc=GREY, lw=Pt(0.75)): return _sp(1,x,y,w,h,fill,lc,lw)
def RND(x,y,w,h, fill=BLUE,  lc=WHITE,lw=Pt(0.75)):
    s=_sp(5,x,y,w,h,fill,lc,lw); s.adjustments[0]=0.09; return s
def DMD(x,y,w,h, fill=BLUE,  lc=WHITE): return _sp(4,x,y,w,h,fill,lc,Pt(0.75))
def OVL(x,y,w,h, fill=NAVY,  lc=WHITE): return _sp(9,x,y,w,h,fill,lc,Pt(0.75))

def T(sp, text, sz=Pt(7), col=WHITE, bold=False, align=PP_ALIGN.CENTER):
    tf=sp.text_frame; p=tf.paragraphs[0]; p.alignment=align
    r=p.runs[0] if p.runs else p.add_run()
    r.text=text; r.font.size=sz; r.font.bold=bold; r.font.color.rgb=col

def LBL(x,y,w,h, text, sz=Pt(7), col=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.size=sz; r.font.bold=bold; r.font.color.rgb=col

def ARR(x1,y1,x2,y2, col=DGREY, w=Pt(1.0)):
    c=sl.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=col; c.line.width=w
    ln=c._element.spPr.find(qn('a:ln'))
    if ln is None: ln=etree.SubElement(c._element.spPr,qn('a:ln'))
    he=ln.find(qn('a:headEnd'))
    if he is None: he=etree.SubElement(ln,qn('a:headEnd'))
    he.set('type','arrow'); he.set('w','med'); he.set('len','med')
    te=ln.find(qn('a:tailEnd'))
    if te is None: te=etree.SubElement(ln,qn('a:tailEnd'))
    te.set('type','none')

# ── Layout constants ───────────────────────────────────────────────────────
LLW = 1.1    # lane label width
LX  = 0.15   # left edge
RX  = 13.18  # right edge
CX  = LX+LLW # content starts at x=1.25
BH  = 0.44   # box height

# Swim lane y-ranges
L1Y, L1H = 0.62, 1.50   # Customer
L2Y, L2H = 2.12, 1.08   # Admin Approval
L3Y, L3H = 3.20, 3.82   # TeleResolve System

# Vertical midpoints
L1M = L1Y + L1H/2        # 1.37
L1T = L1M - BH/2         # 1.15
L2M = L2Y + L2H/2        # 2.66

R1M = L3Y + 0.70          # row-1 mid  3.90
R2M = L3Y + 1.95          # row-2 mid  5.15
R3M = L3Y + 3.20          # row-3 mid  6.40
R1T = R1M - BH/2          # 3.68
R2T = R2M - BH/2          # 4.93
R3T = R3M - BH/2          # 6.18

# Row-1: 5 boxes spanning content width
BW1 = 2.0
GAP = 0.14
B1X = CX + 0.18
R1XS = [B1X + i*(BW1+GAP) for i in range(5)]  # left-x of each box

# Row-2 boxes
SIC_X, SIC_W = 4.70, 1.65
VD_X,  VD_W  = SIC_X+SIC_W+0.18, 1.65
VD_CX = VD_X + VD_W/2

# Row-3 items (no X circles, just direct arrows)
SAC_X, SAC_W = CX+0.18, 1.0          # Start AI Chat oval
RND_X, RND_W = SAC_X+SAC_W+0.20, 1.55  # Run Network Diagnosis
LD_X,  LD_W  = RND_X+RND_W+0.20, 1.50  # Live Dashboard
END_X, END_W = LD_X+LD_W+0.20, 0.75    # End oval

# ── Slide chrome ───────────────────────────────────────────────────────────
RCT(0, 0, 13.33, 7.5, fill=WHITE, lc=None)

# Header
RCT(0, 0, 13.33, 0.52, fill=NAVY, lc=None)
LBL(0.18, 0.06, 6, 0.22, "6. Business Architecture", sz=Pt(13), col=WHITE, bold=True)
LBL(0.18, 0.30, 10, 0.18,
    "Business Process (To-Be)  —  TeleResolve AI-Powered Complaint Handling System",
    sz=Pt(8.5), col=CYAN)

# Footer
RCT(0, 7.12, 13.33, 0.38, fill=NAVY, lc=None)
LBL(0.2, 7.14, 12.9, 0.22,
    "© 2026 KPMG Assurance and Consulting Services LLP an Indian Limited Liability Partnership and a member firm of the KPMG global organization of independent member firms affiliated with KPMG International Limited, a private English company limited by guarantee. All rights reserved.",
    sz=Pt(6), col=GREY)

# Diagram outer border
RCT(LX, 0.60, RX-LX, L3Y+L3H-0.60, fill=WHITE, lc=GREY, lw=Pt(1.0))

# ── Swim lane bands (WHITE background, no colour) ─────────────────────────
for name, ly, lh in [
    ("Customer",            L1Y, L1H),
    ("Admin\nApproval",     L2Y, L2H),
    ("TeleResolve\nSystem", L3Y, L3H),
]:
    sp = RCT(LX, ly, LLW, lh, fill=NAVY, lc=WHITE, lw=Pt(0.5))
    T(sp, name, sz=Pt(8.5), col=WHITE, bold=True)
    RCT(CX, ly, RX-CX, lh, fill=WHITE, lc=GREY, lw=Pt(0.5))  # white content bg

# ═══════════════════════════════════════════════════════════════════════════
# LANE 1 — CUSTOMER
# ═══════════════════════════════════════════════════════════════════════════
# Start oval
sp = OVL(CX+0.05, L1T, 0.72, BH); T(sp, "Start", Pt(8), WHITE, bold=True)
S_CX = CX+0.05+0.36
ARR(S_CX+0.36, L1M, S_CX+0.54, L1M)

# Diamond: Existing User?
DW, DH = 1.05, 0.68
DX = S_CX+0.54
D_CX  = DX+DW/2
D_TOP = L1T-(DH-BH)/2
sp = DMD(DX, D_TOP, DW, DH); T(sp, "Existing\nUser?", Pt(6.5), WHITE)

# YES arc — over the top of lane 1 → Login
ARC_Y = L1Y+0.07
LBL(D_CX+0.04, ARC_Y-0.02, 0.28, 0.18, "Yes", sz=Pt(6.5), col=DGREY)
ARR(D_CX, D_TOP,  D_CX,  ARC_Y)
ARR(D_CX, ARC_Y,  12.30, ARC_Y)
ARR(12.30, ARC_Y, 12.30, L1T)

# NO → Register
NO_X = DX+DW
LBL(NO_X+0.04, L1M-0.14, 0.25, 0.18, "No", sz=Pt(6.5), col=DGREY)
ARR(NO_X, L1M, NO_X+0.22, L1M)

REG_X, REG_W = NO_X+0.22, 1.05
sp = RND(REG_X, L1T, REG_W, BH); T(sp, "Register", Pt(7), WHITE)
REG_CX  = REG_X+REG_W/2
REG_BOT = L1T+BH
ARR(REG_X+REG_W, L1M, REG_X+REG_W+0.18, L1M)

CP_X, CP_W = REG_X+REG_W+0.18, 1.20
sp = RND(CP_X, L1T, CP_W, BH); T(sp, "Create\nPassword", Pt(6.5), WHITE)
ARR(CP_X+CP_W, L1M, CP_X+CP_W+0.18, L1M)

OTP_X, OTP_W = CP_X+CP_W+0.18, 1.20
sp = RND(OTP_X, L1T, OTP_W, BH); T(sp, "Generate\nOTP", Pt(6.5), WHITE)
ARR(OTP_X+OTP_W, L1M, 11.75, L1M)

# Login
LOGIN_X, LOGIN_W = 11.75, 1.0
LOGIN_CX = LOGIN_X+LOGIN_W/2
LOGIN_BOT = L1T+BH
sp = RND(LOGIN_X, L1T, LOGIN_W, BH, fill=NAVY); T(sp, "Login", Pt(8.5), WHITE, bold=True)

# ═══════════════════════════════════════════════════════════════════════════
# LANE 2 — ADMIN APPROVAL
# ═══════════════════════════════════════════════════════════════════════════
ADW, ADH = 1.05, 0.70
AD_CX = REG_CX
ADX   = AD_CX-ADW/2
AD_TOP= L2M-ADH/2
sp = DMD(ADX, AD_TOP, ADW, ADH); T(sp, "Admin\nApproval", Pt(6.5), WHITE)

ARR(REG_CX, REG_BOT, REG_CX, AD_TOP)            # Register → Admin diamond

LBL(ADX-0.30, L2M-0.10, 0.25, 0.18, "No",  sz=Pt(6.5), col=DGREY)
ARR(ADX, L2M, CX+0.10, L2M)
sp = RND(CX+0.10, L2M-BH/2, 0.95, BH, fill=RED); T(sp, "Reject &\nNotify", Pt(6.5), WHITE)

LBL(ADX+ADW+0.04, L2M-0.10, 0.28, 0.18, "Yes", sz=Pt(6.5), col=DGREY)
ARR(ADX+ADW, L2M, LOGIN_CX, L2M)                # Yes → right to Login col
ARR(LOGIN_CX, L2M, LOGIN_CX, LOGIN_BOT)          # up to Login bottom

# ═══════════════════════════════════════════════════════════════════════════
# LANE 3 — TELERESOLVE SYSTEM  (no X circles)
# ═══════════════════════════════════════════════════════════════════════════

# Login → drop into lane 3 → L-shape into Row-1 box 1
ARR(LOGIN_CX, LOGIN_BOT, LOGIN_CX, L3Y+0.22)    # down
ARR(LOGIN_CX, L3Y+0.22,  B1X,      L3Y+0.22)    # left along top of lane 3
ARR(B1X,      L3Y+0.22,  B1X,      R1T)          # down into row-1 box 1 top

# ── ROW 1: 5 feature boxes ────────────────────────────────────────────────
ROW1 = [
    "Submit\nComplaint",
    "View Active\nTickets",
    "Recent\nSessions",
    "Network AI\nChat",
    "Upload\nEvidence",
]
for i, (rx0, txt) in enumerate(zip(R1XS, ROW1)):
    sp = RND(rx0, R1T, BW1, BH); T(sp, txt, Pt(6.5), WHITE)
    if i < 4:
        ARR(rx0+BW1, R1M, R1XS[i+1], R1M)

# Row-1 last box right edge → down-right → Row-2 (Select Issue Category)
LAST_R = R1XS[4]+BW1
ARR(LAST_R,  R1M,   LAST_R+0.25, R1M)
ARR(LAST_R+0.25, R1M, LAST_R+0.25, R2M)          # down
ARR(LAST_R+0.25, R2M, SIC_X+SIC_W, R2M)          # left to Select right edge

# ── ROW 2: Select Issue Category → View Dashboard ─────────────────────────
sp = RND(SIC_X, R2T, SIC_W, BH); T(sp, "Select Issue\nCategory", Pt(6.5), WHITE)
ARR(SIC_X+SIC_W, R2M, VD_X, R2M)
sp = RND(VD_X,  R2T, VD_W,  BH); T(sp, "View\nDashboard", Pt(6.5), WHITE)

# View Dashboard → down → into Row-3 (between Start AI Chat and Run Diagnosis)
ARR(VD_CX, R2T+BH, VD_CX, R3M)                   # straight down
ARR(VD_CX, R3M, SAC_X+SAC_W+0.10, R3M)           # left to Start AI Chat right area

# ── ROW 3: Start AI Chat → Run Network Diagnosis → Live Dashboard → End ───
# Arrow from lane-3 left edge down into Start AI Chat
ARR(B1X, R1T, B1X, R3M)                           # left column drop (vertical guide)
ARR(B1X, R3M, SAC_X, R3M)                         # right into Start AI Chat

sp = OVL(SAC_X, R3T, SAC_W, BH, fill=NAVY); T(sp, "Start AI\nChat", Pt(6.5), WHITE)
ARR(SAC_X+SAC_W, R3M, RND_X, R3M)

sp = RND(RND_X, R3T, RND_W, BH); T(sp, "Run Network\nDiagnosis", Pt(6.5), WHITE)
ARR(RND_X+RND_W, R3M, LD_X, R3M)

sp = RND(LD_X, R3T, LD_W, BH); T(sp, "Live\nDashboard", Pt(6.5), WHITE)
ARR(LD_X+LD_W, R3M, END_X, R3M)

sp = OVL(END_X, R3T, END_W, BH, fill=NAVY); T(sp, "End", Pt(8.5), WHITE, bold=True)

# ── Save ───────────────────────────────────────────────────────────────────
OUT = r"c:\Users\didar\Downloads\files (12)\telecom-complaint-system\TeleResolve_ToBe_Clean.pptx"
prs.save(OUT)
print("Saved: " + OUT)
