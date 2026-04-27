"""
TeleResolve - Business Architecture PowerPoint Generator
Generates slides matching KPMG-style diagrams adapted to the TeleResolve project
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Color Palette (matching KPMG navy/blue style) ──────────────────────────
KPMG_NAVY   = RGBColor(0x00, 0x33, 0x6A)   # dark navy
KPMG_BLUE   = RGBColor(0x00, 0x5B, 0xB5)   # medium blue
KPMG_LBLUE  = RGBColor(0x41, 0x72, 0xC4)   # light blue (swim-lane boxes)
KPMG_CYAN   = RGBColor(0x00, 0xB0, 0xF0)   # accent cyan
KPMG_GREEN  = RGBColor(0x70, 0xAD, 0x47)   # arrow green (roadmap)
KPMG_AMBER  = RGBColor(0xED, 0x7D, 0x31)   # platform box amber
KPMG_GREY   = RGBColor(0xF2, 0xF2, 0xF2)   # slide background
KPMG_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
KPMG_BLACK  = RGBColor(0x00, 0x00, 0x00)
KPMG_RED    = RGBColor(0xFF, 0x00, 0x00)
SWIMLANE_BG = RGBColor(0xDA, 0xE3, 0xF3)   # very light blue

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Helper utilities ────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1.2), radius=0.05):
    from pptx.util import Emu
    sp = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sp.adjustments[0] = radius
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = line
        sp.line.width = line_w
    else:
        sp.line.fill.background()
    return sp


def add_diamond(slide, x, y, w, h, fill=KPMG_BLUE, line=KPMG_WHITE):
    sp = slide.shapes.add_shape(
        4,  # DIAMOND
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line
    sp.line.width = Pt(1)
    return sp


def add_oval(slide, x, y, w, h, fill=KPMG_NAVY, line=KPMG_WHITE):
    sp = slide.shapes.add_shape(
        9,  # OVAL
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line
    sp.line.width = Pt(1)
    return sp


def set_text(shape, text, size=Pt(7), bold=False, color=KPMG_WHITE, align=PP_ALIGN.CENTER, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def add_label(slide, x, y, w, h, text, size=Pt(7), bold=False, color=KPMG_WHITE, align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_arrow_line(slide, x1, y1, x2, y2, color=KPMG_NAVY, w=Pt(1.5)):
    """Add a line with arrow end between two points (in inches)."""
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = w
    # Add arrowhead via XML
    ln = connector._element.spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(connector._element.spPr, qn('a:ln'))
    tail_e = etree.SubElement(ln, qn('a:tailEnd'))
    tail_e.set('type', 'none')
    head_e = ln.find(qn('a:headEnd'))
    if head_e is None:
        head_e = etree.SubElement(ln, qn('a:headEnd'))
    head_e.set('type', 'arrow')
    head_e.set('w', 'med')
    head_e.set('len', 'med')


def box_with_text(slide, x, y, w, h, text, fill=KPMG_LBLUE, line=KPMG_WHITE,
                  tsize=Pt(6.5), bold=False, tcolor=KPMG_WHITE, shape_type='rect'):
    if shape_type == 'oval':
        sp = add_oval(slide, x, y, w, h, fill=fill, line=line)
    elif shape_type == 'diamond':
        sp = add_diamond(slide, x, y, w, h, fill=fill, line=line)
    elif shape_type == 'rounded':
        sp = add_rounded_rect(slide, x, y, w, h, fill=fill, line=line)
    else:
        sp = add_rect(slide, x, y, w, h, fill=fill, line=line)
    set_text(sp, text, size=tsize, bold=bold, color=tcolor)
    return sp


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════
blank_layout = prs.slide_layouts[6]  # completely blank

slide1 = prs.slides.add_slide(blank_layout)

# Background
bg = add_rect(slide1, 0, 0, 13.33, 7.5, fill=KPMG_NAVY)

# Top accent bar
add_rect(slide1, 0, 0, 13.33, 0.08, fill=KPMG_CYAN)

# Logo placeholder (text)
add_label(slide1, 0.3, 0.15, 1.5, 0.4, "KPMG", size=Pt(18), bold=True, color=KPMG_WHITE)

# Title
t = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1.2))
tf = t.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "TeleResolve — Telecom Complaint Handling System"
r.font.size = Pt(30)
r.font.bold = True
r.font.color.rgb = KPMG_WHITE

# Subtitle
t2 = slide1.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9), Inches(0.6))
tf2 = t2.text_frame
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Business Architecture — As-Is Process | To-Be Process | Roadmap"
r2.font.size = Pt(14)
r2.font.color.rgb = KPMG_CYAN

# Divider line
add_rect(slide1, 1.5, 4.2, 10.3, 0.03, fill=KPMG_CYAN)

# Meta info
add_label(slide1, 1.5, 4.4, 10, 0.3,
          "Prepared by: TeleResolve Team  |  April 2026  |  Confidential",
          size=Pt(10), color=KPMG_LBLUE, align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide1, 0, 7.1, 13.33, 0.4, fill=RGBColor(0x00, 0x1F, 0x4D))
add_label(slide1, 0.2, 7.15, 12, 0.25,
          "© 2026 TeleResolve. All rights reserved.",
          size=Pt(8), color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — BUSINESS PROCESS (AS-IS)
# ═══════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)

# Slide background
add_rect(slide2, 0, 0, 13.33, 7.5, fill=KPMG_GREY)

# Header bar
add_rect(slide2, 0, 0, 13.33, 0.55, fill=KPMG_NAVY)
add_label(slide2, 0.2, 0.08, 4, 0.4, "6. Business Architecture", size=Pt(14), bold=True, color=KPMG_WHITE)
add_label(slide2, 0.2, 0.35, 5, 0.25, "Business Process (As-Is) — Manual Complaint Handling", size=Pt(9), color=KPMG_CYAN)

# Diagram outer border
add_rect(slide2, 0.15, 0.65, 12.9, 6.4, fill=KPMG_WHITE, line=KPMG_NAVY, line_w=Pt(1.5))

# ── Swim Lane Labels (left column) ─────────────────────────────────────────
lane_x = 0.2
lane_label_w = 1.1
lane_h = 1.5

lanes = [
    ("Customer", 0.7),
    ("Contact\nCenter Agent", 2.25),
    ("Manager /\nSupervisor", 3.8),
    ("Network\nEngineer", 5.35),
]

for label, ly in lanes:
    add_rect(slide2, lane_x, ly, lane_label_w, lane_h,
             fill=KPMG_NAVY, line=KPMG_WHITE, line_w=Pt(0.8))
    add_label(slide2, lane_x + 0.02, ly + 0.3, lane_label_w - 0.05, lane_h - 0.6,
              label, size=Pt(8), bold=True, color=KPMG_WHITE, align=PP_ALIGN.CENTER)
    # horizontal separator
    add_rect(slide2, lane_x + lane_label_w, ly + lane_h, 12.9 - lane_label_w - 0.1, 0.02,
             fill=RGBColor(0xCC, 0xCC, 0xCC))

# Light swim lane BG rows
row_colors = [
    RGBColor(0xE9, 0xF0, 0xFB),
    RGBColor(0xF5, 0xF8, 0xFF),
    RGBColor(0xE9, 0xF0, 0xFB),
    RGBColor(0xF5, 0xF8, 0xFF),
]
for i, (_, ly) in enumerate(lanes):
    add_rect(slide2, lane_x + lane_label_w, ly, 12.9 - lane_label_w - 0.1, lane_h,
             fill=row_colors[i], line=None)

# ── CUSTOMER LANE — Row 1 (y=0.7) ──────────────────────────────────────────
row1_y = 0.82
# Start oval
box_with_text(slide2, 1.45, row1_y, 0.65, 0.45, "Start", fill=KPMG_NAVY, shape_type='oval', tsize=Pt(7.5), bold=True)
add_arrow_line(slide2, 2.1, row1_y + 0.225, 2.5, row1_y + 0.225)
# Existing Customer?
box_with_text(slide2, 2.5, row1_y - 0.08, 0.85, 0.62, "Existing\nCustomer?", fill=KPMG_BLUE, shape_type='diamond', tsize=Pt(6.5))
# No → Register (Manual Form)
add_label(slide2, 3.37, row1_y + 0.05, 0.25, 0.2, "No", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide2, 3.35, row1_y + 0.23, 3.75, row1_y + 0.23)
box_with_text(slide2, 3.75, row1_y, 0.85, 0.45, "Register\n(Manual Form)", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 4.6, row1_y + 0.225, 5.0, row1_y + 0.225)
box_with_text(slide2, 5.0, row1_y, 0.9, 0.45, "Phone / Web\nComplaint", fill=KPMG_LBLUE, tsize=Pt(6.5))
# Yes arrow down to phone complaint
add_label(slide2, 2.85, row1_y + 0.62, 0.3, 0.2, "Yes", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide2, 2.93, row1_y + 0.62, 2.93, row1_y + 1.5)  # down to row2

# ── CONTACT CENTER AGENT LANE — Row 2 (y=2.25) ─────────────────────────────
row2_y = 2.38
box_with_text(slide2, 1.45, row2_y, 0.9, 0.45, "Receive\nCall/Request", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 2.35, row2_y + 0.225, 2.75, row2_y + 0.225)
box_with_text(slide2, 2.75, row2_y, 0.9, 0.45, "Manual\nCategorization", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 3.65, row2_y + 0.225, 4.05, row2_y + 0.225)
box_with_text(slide2, 4.05, row2_y, 0.9, 0.45, "Create Ticket\n(Spreadsheet)", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 4.95, row2_y + 0.225, 5.35, row2_y + 0.225)
box_with_text(slide2, 5.35, row2_y - 0.08, 0.85, 0.62, "Resolved\nby Agent?", fill=KPMG_BLUE, shape_type='diamond', tsize=Pt(6.5))
add_label(slide2, 6.22, row2_y + 0.05, 0.3, 0.2, "Yes", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide2, 6.2, row2_y + 0.23, 6.6, row2_y + 0.23)
box_with_text(slide2, 6.6, row2_y, 0.9, 0.45, "Notify\nCustomer", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 7.5, row2_y + 0.225, 7.9, row2_y + 0.225)
box_with_text(slide2, 7.9, row2_y, 0.9, 0.45, "Close &\nManual Report", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 8.8, row2_y + 0.225, 9.2, row2_y + 0.225)
box_with_text(slide2, 9.2, row2_y, 0.6, 0.45, "End", fill=KPMG_NAVY, shape_type='oval', tsize=Pt(7.5), bold=True)
# No escalate to Manager
add_label(slide2, 5.75, row2_y + 0.62, 0.25, 0.2, "No", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide2, 5.78, row2_y + 0.62, 5.78, row2_y + 1.5)

# ── MANAGER LANE — Row 3 (y=3.8) ───────────────────────────────────────────
row3_y = 3.93
box_with_text(slide2, 1.45, row3_y, 0.9, 0.45, "Manual\nEscalation", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 2.35, row3_y + 0.225, 2.75, row3_y + 0.225)
box_with_text(slide2, 2.75, row3_y, 0.9, 0.45, "Review &\nPrioritize", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 3.65, row3_y + 0.225, 4.05, row3_y + 0.225)
box_with_text(slide2, 4.05, row3_y - 0.08, 0.85, 0.62, "Network\nIssue?", fill=KPMG_BLUE, shape_type='diamond', tsize=Pt(6.5))
add_label(slide2, 4.92, row3_y + 0.05, 0.3, 0.2, "No", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide2, 4.9, row3_y + 0.23, 5.3, row3_y + 0.23)
box_with_text(slide2, 5.3, row3_y, 0.9, 0.45, "Manual\nResolution", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 6.2, row3_y + 0.225, 7.5, row2_y + 0.45)  # back to notify customer
add_label(slide2, 4.47, row3_y + 0.62, 0.3, 0.2, "Yes", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide2, 4.48, row3_y + 0.62, 4.48, row3_y + 1.5)

# ── NETWORK ENGINEER LANE — Row 4 (y=5.35) ─────────────────────────────────
row4_y = 5.48
box_with_text(slide2, 1.45, row4_y, 0.9, 0.45, "Manual Site\nInvestigation", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 2.35, row4_y + 0.225, 2.75, row4_y + 0.225)
box_with_text(slide2, 2.75, row4_y, 0.9, 0.45, "Identify Root\nCause (Manual)", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 3.65, row4_y + 0.225, 4.05, row4_y + 0.225)
box_with_text(slide2, 4.05, row4_y, 0.9, 0.45, "Parameter\nChange (Manual)", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 4.95, row4_y + 0.225, 5.35, row4_y + 0.225)
box_with_text(slide2, 5.35, row4_y, 1.0, 0.45, "Email Approval\n(No Tracking)", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 6.35, row4_y + 0.225, 6.75, row4_y + 0.225)
box_with_text(slide2, 6.75, row4_y, 0.9, 0.45, "Implement\nChange", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide2, 7.65, row4_y + 0.225, 5.78, row3_y + 0.45)  # up to manager lane

# Pain points box
pain = add_rounded_rect(slide2, 9.5, 2.0, 3.3, 4.5,
                        fill=RGBColor(0xFF, 0xF2, 0xCC), line=KPMG_AMBER, line_w=Pt(1.5))
pain.fill.solid()
pain.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
pain.line.color.rgb = KPMG_AMBER
add_label(slide2, 9.6, 2.05, 3.1, 0.35, "Pain Points (As-Is)",
          size=Pt(9), bold=True, color=KPMG_AMBER, align=PP_ALIGN.LEFT)
pain_points = [
    "• No unified complaint portal",
    "• Manual ticket creation in spreadsheets",
    "• No AI/ML-assisted diagnosis",
    "• Slow escalation (phone/email)",
    "• No real-time SLA tracking",
    "• No automated network diagnostics",
    "• No change request audit trail",
    "• Customer has no visibility on status",
    "• No CSAT / feedback collection",
    "• No analytics or KPI dashboard",
]
for i, pt in enumerate(pain_points):
    add_label(slide2, 9.6, 2.45 + i * 0.37, 3.1, 0.35,
              pt, size=Pt(7.5), color=KPMG_BLACK, align=PP_ALIGN.LEFT)

# Footer
add_rect(slide2, 0, 7.1, 13.33, 0.4, fill=KPMG_NAVY)
add_label(slide2, 0.2, 7.15, 12, 0.25,
          "© 2026 TeleResolve. Confidential  |  Business Architecture — As-Is Process",
          size=Pt(8), color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BUSINESS PROCESS (TO-BE)
# ═══════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
add_rect(slide3, 0, 0, 13.33, 7.5, fill=KPMG_GREY)

# Header
add_rect(slide3, 0, 0, 13.33, 0.55, fill=KPMG_NAVY)
add_label(slide3, 0.2, 0.08, 4, 0.4, "6. Business Architecture", size=Pt(14), bold=True, color=KPMG_WHITE)
add_label(slide3, 0.2, 0.35, 6, 0.25, "Business Process (To-Be) — TeleResolve AI-Powered Complaint System", size=Pt(9), color=KPMG_CYAN)

# Outer border
add_rect(slide3, 0.15, 0.65, 12.9, 6.4, fill=KPMG_WHITE, line=KPMG_NAVY, line_w=Pt(1.5))

# ── Swim lane definitions ───────────────────────────────────────────────────
s3_lanes = [
    ("Customer",           0.70,  1.45),
    ("Admin\nApproval",    2.20,  1.10),
    ("TeleResolve\nSystem",3.35,  1.55),
    ("Agent /\nEngineer",  4.95,  1.55),
]

for label, ly, lh in s3_lanes:
    add_rect(slide3, 0.2, ly, 1.1, lh, fill=KPMG_NAVY, line=KPMG_WHITE, line_w=Pt(0.8))
    add_label(slide3, 0.22, ly + 0.28, 1.06, lh - 0.56,
              label, size=Pt(8), bold=True, color=KPMG_WHITE, align=PP_ALIGN.CENTER)
    bg_col = RGBColor(0xE9, 0xF0, 0xFB) if s3_lanes.index((label, ly, lh)) % 2 == 0 else RGBColor(0xF5, 0xF8, 0xFF)
    add_rect(slide3, 1.3, ly, 11.7, lh, fill=bg_col, line=None)
    add_rect(slide3, 0.2, ly + lh, 12.8, 0.02, fill=RGBColor(0xCC, 0xCC, 0xCC))

# ── ROW 1: CUSTOMER LANE (y=0.70–2.15) ──────────────────────────────────────
r1y = 0.85
box_with_text(slide3, 1.4,  r1y, 0.6,  0.42, "Start", fill=KPMG_NAVY, shape_type='oval', tsize=Pt(7.5), bold=True)
add_arrow_line(slide3, 2.0, r1y+0.21, 2.35, r1y+0.21)
box_with_text(slide3, 2.35, r1y-0.07, 0.82, 0.58, "Existing\nUser?", fill=KPMG_BLUE, shape_type='diamond', tsize=Pt(6.5))
# No → Register
add_label(slide3, 3.19, r1y+0.1, 0.25, 0.2, "No", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide3, 3.17, r1y+0.21, 3.55, r1y+0.21)
box_with_text(slide3, 3.55, r1y, 0.88, 0.42, "Self-Register\n(Web/App)", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 4.43, r1y+0.21, 4.83, r1y+0.21)
box_with_text(slide3, 4.83, r1y, 0.88, 0.42, "Create\nPassword", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 5.71, r1y+0.21, 6.11, r1y+0.21)
box_with_text(slide3, 6.11, r1y, 0.88, 0.42, "OTP\nVerification", fill=KPMG_LBLUE, tsize=Pt(6.5))
# Yes arrow goes to Login
add_label(slide3, 2.72, r1y+0.58, 0.3, 0.2, "Yes", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide3, 2.76, r1y+0.58, 2.76, r1y+1.02)  # arc to Login below (we draw below)
# Login box
box_with_text(slide3, 10.6, r1y, 0.88, 0.42, "Login to\nTeleResolve", fill=KPMG_NAVY, tsize=Pt(6.5))
add_arrow_line(slide3, 6.99, r1y+0.21, 10.6, r1y+0.21)  # from OTP to Login
# Yes shortcut line from diamond direct to Login
add_arrow_line(slide3, 11.04, r1y+0.21, 11.48, r1y+0.21)
box_with_text(slide3, 11.48, r1y, 0.88, 0.42, "Select\nSector", fill=KPMG_LBLUE, tsize=Pt(6.5))
# Down from Login to TeleResolve lane
add_arrow_line(slide3, 11.92, r1y+0.42, 11.92, 3.5)

# ── ROW 2: ADMIN APPROVAL LANE (y=2.20–3.30) ────────────────────────────────
r2y = 2.33
# Admin approval diamond (for new registrations)
add_arrow_line(slide3, 3.99, r1y+0.21, 3.99, r2y+0.07)  # down from Register
box_with_text(slide3, 3.55, r2y, 0.88, 0.58, "Admin\nApproval?", fill=KPMG_BLUE, shape_type='diamond', tsize=Pt(6.5))
add_label(slide3, 4.45, r2y+0.14, 0.3, 0.2, "Yes", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide3, 4.43, r2y+0.29, 4.83, r2y+0.29)
box_with_text(slide3, 4.83, r2y+0.05, 0.88, 0.42, "Account\nActivated", fill=KPMG_GREEN, tsize=Pt(6.5))
add_arrow_line(slide3, 5.71, r2y+0.26, 6.11, r2y+0.26)
box_with_text(slide3, 6.11, r2y+0.05, 0.88, 0.42, "Email OTP\nSent", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 6.99, r2y+0.26, 4.83, r1y+0.21)  # back to Create Password row1
# No arrow
add_label(slide3, 3.1, r2y+0.36, 0.3, 0.2, "No", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide3, 3.12, r2y+0.38, 2.3, r2y+0.38)
box_with_text(slide3, 1.4, r2y+0.15, 0.88, 0.42, "Rejection\nEmail Sent", fill=KPMG_RED, tsize=Pt(6.5))

# ── ROW 3: TELERESOLVE SYSTEM LANE (y=3.35–4.90) ────────────────────────────
r3y = 3.5
box_with_text(slide3, 1.4,  r3y, 0.92, 0.42, "AI Chatbot\nSession", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 2.32, r3y+0.21, 2.72, r3y+0.21)
box_with_text(slide3, 2.72, r3y-0.07, 0.82, 0.58, "Resolved\nby AI?", fill=KPMG_BLUE, shape_type='diamond', tsize=Pt(6.5))
# Yes → Feedback
add_label(slide3, 3.56, r3y+0.1, 0.3, 0.2, "Yes", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide3, 3.54, r3y+0.21, 3.94, r3y+0.21)
box_with_text(slide3, 3.94, r3y, 0.88, 0.42, "Customer\nFeedback", fill=KPMG_GREEN, tsize=Pt(6.5))
add_arrow_line(slide3, 4.82, r3y+0.21, 5.22, r3y+0.21)
box_with_text(slide3, 5.22, r3y, 0.88, 0.42, "Session\nClosed", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 6.1, r3y+0.21, 6.5, r3y+0.21)
box_with_text(slide3, 6.5, r3y, 0.55, 0.42, "End", fill=KPMG_NAVY, shape_type='oval', tsize=Pt(7.5), bold=True)
# No → Escalate
add_label(slide3, 3.1, r3y+0.38, 0.3, 0.2, "No", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide3, 3.13, r3y+0.38, 3.13, r3y+0.72)
box_with_text(slide3, 2.72, r3y+0.72, 0.82, 0.42, "Escalate\nto Ticket", fill=KPMG_AMBER, tsize=Pt(6.5))
add_arrow_line(slide3, 3.54, r3y+0.93, 3.94, r3y+0.93)
box_with_text(slide3, 3.94, r3y+0.72, 0.88, 0.42, "SLA Timer\nStarted", fill=KPMG_AMBER, tsize=Pt(6.5))
add_arrow_line(slide3, 4.82, r3y+0.93, 4.82, r3y+1.5)  # down to agent lane

# Live Dashboard boxes
box_with_text(slide3, 8.0, r3y, 0.9, 0.42, "Live KPI\nDashboard", fill=KPMG_LBLUE, tsize=Pt(6.5))
box_with_text(slide3, 9.1, r3y, 0.9, 0.42, "SLA Alert\nMonitor", fill=KPMG_LBLUE, tsize=Pt(6.5))
box_with_text(slide3, 10.2, r3y, 0.9, 0.42, "Analytics &\nReports", fill=KPMG_LBLUE, tsize=Pt(6.5))
box_with_text(slide3, 11.3, r3y, 0.9, 0.42, "Change\nRequests", fill=KPMG_LBLUE, tsize=Pt(6.5))

# ── ROW 4: AGENT/ENGINEER LANE (y=4.95–6.50) ────────────────────────────────
r4y = 5.1
box_with_text(slide3, 1.4,  r4y, 0.92, 0.42, "Ticket\nAssigned", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 2.32, r4y+0.21, 2.72, r4y+0.21)
box_with_text(slide3, 2.72, r4y, 0.88, 0.42, "Network AI\nDiagnosis", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 3.6, r4y+0.21, 4.0, r4y+0.21)
box_with_text(slide3, 4.0, r4y-0.07, 0.85, 0.58, "Network\nIssue?", fill=KPMG_BLUE, shape_type='diamond', tsize=Pt(6.5))
# Yes
add_label(slide3, 4.87, r4y+0.1, 0.3, 0.2, "Yes", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide3, 4.85, r4y+0.21, 5.25, r4y+0.21)
box_with_text(slide3, 5.25, r4y, 0.92, 0.42, "Root Cause\nAnalysis (AI)", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 6.17, r4y+0.21, 6.57, r4y+0.21)
box_with_text(slide3, 6.57, r4y, 0.92, 0.42, "CR / Param\nChange", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 7.49, r4y+0.21, 7.89, r4y+0.21)
box_with_text(slide3, 7.89, r4y, 0.92, 0.42, "Mgr / CTO\nApproval", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 8.81, r4y+0.21, 9.21, r4y+0.21)
box_with_text(slide3, 9.21, r4y, 0.92, 0.42, "Implement\nChange", fill=KPMG_LBLUE, tsize=Pt(6.5))
# No
add_label(slide3, 4.4, r4y+0.36, 0.3, 0.2, "No", size=Pt(6.5), color=KPMG_NAVY)
add_arrow_line(slide3, 4.43, r4y+0.38, 4.43, r4y+0.82)
box_with_text(slide3, 4.0, r4y+0.82, 0.88, 0.42, "Billing /\nSIM Issue Fix", fill=KPMG_LBLUE, tsize=Pt(6.5))
add_arrow_line(slide3, 4.88, r4y+1.03, 9.21, r4y+0.42)  # merge to Implement Change

# Ticket resolved → feedback
add_arrow_line(slide3, 10.13, r4y+0.21, 10.53, r4y+0.21)
box_with_text(slide3, 10.53, r4y, 0.92, 0.42, "Ticket\nResolved", fill=KPMG_GREEN, tsize=Pt(6.5))
add_arrow_line(slide3, 11.45, r4y+0.21, 11.85, r4y+0.21)
box_with_text(slide3, 11.85, r4y, 0.6, 0.42, "End", fill=KPMG_NAVY, shape_type='oval', tsize=Pt(7.5), bold=True)

# Footer
add_rect(slide3, 0, 7.1, 13.33, 0.4, fill=KPMG_NAVY)
add_label(slide3, 0.2, 7.15, 12, 0.25,
          "© 2026 TeleResolve. Confidential  |  Business Architecture — To-Be Process",
          size=Pt(8), color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ROADMAP
# ═══════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
add_rect(slide4, 0, 0, 13.33, 7.5, fill=KPMG_GREY)

# Header
add_rect(slide4, 0, 0, 13.33, 0.55, fill=KPMG_NAVY)
add_label(slide4, 0.2, 0.08, 4, 0.4, "7. Roadmap", size=Pt(14), bold=True, color=KPMG_WHITE)
add_label(slide4, 0.2, 0.35, 8, 0.25, "TeleResolve Platform Implementation Roadmap", size=Pt(9), color=KPMG_CYAN)

# Outer border
add_rect(slide4, 0.15, 0.65, 12.9, 6.4, fill=KPMG_WHITE, line=KPMG_NAVY, line_w=Pt(1.5))

# Platform label box (amber like KPMG diagram)
plat = add_rounded_rect(slide4, 0.4, 1.2, 1.7, 1.0, fill=KPMG_AMBER, line=KPMG_WHITE)
add_label(slide4, 0.42, 1.4, 1.66, 0.6, "TeleResolve\nPlatform", size=Pt(10), bold=True, color=KPMG_WHITE, align=PP_ALIGN.CENTER)

# Arrow phases — use chevron-style rectangles
def add_phase_arrow(slide, x, y, w, h, label, fill):
    """Draw a simple arrow/chevron shape representing a phase."""
    sp = add_rounded_rect(slide, x, y, w, h, fill=fill, line=KPMG_WHITE, line_w=Pt(2), radius=0.02)
    add_label(slide, x + 0.08, y + 0.1, w - 0.15, h - 0.2, label,
              size=Pt(11), bold=True, color=KPMG_WHITE, align=PP_ALIGN.CENTER)
    return sp

phase_y = 1.05
phase_h = 0.55
phase_w = 2.8

add_phase_arrow(slide4, 2.4,  phase_y, phase_w, phase_h, "Current State",    KPMG_GREEN)
add_label(slide4, 5.15, phase_y + 0.15, 0.5, 0.3, "→", size=Pt(20), bold=True, color=KPMG_GREEN)
add_phase_arrow(slide4, 5.6,  phase_y, phase_w, phase_h, "Transition State", KPMG_GREEN)
add_label(slide4, 8.35, phase_y + 0.15, 0.5, 0.3, "→", size=Pt(20), bold=True, color=KPMG_GREEN)
add_phase_arrow(slide4, 8.8,  phase_y, phase_w, phase_h, "Target State",     KPMG_GREEN)

# ── Current State Content ─────────────────────────────────────────────────
cs_x = 2.4
current_items = [
    ("Phase-In (Start):", "Pilot deployment to 1 telecom zone; Core chatbot, ticket & agent module live"),
    ("Phase-In (End):",   "Pilot readiness validated; Customer & Agent roles tested"),
    ("Phase-Out (Start):","Manual spreadsheet complaints begin reduction"),
    ("Phase-Out (End):",  "Pilot zone fully on TeleResolve; manual process retired"),
]
add_rect(slide4, cs_x, 1.7, phase_w, 4.8, fill=RGBColor(0xF0, 0xF5, 0xFF), line=KPMG_LBLUE, line_w=Pt(0.8))
for i, (title, desc) in enumerate(current_items):
    ty = 1.8 + i * 1.15
    add_label(slide4, cs_x + 0.1, ty, phase_w - 0.2, 0.28, title,
              size=Pt(8.5), bold=True, color=KPMG_NAVY, align=PP_ALIGN.LEFT)
    add_label(slide4, cs_x + 0.1, ty + 0.27, phase_w - 0.2, 0.75, desc,
              size=Pt(7.5), color=KPMG_BLACK, align=PP_ALIGN.LEFT)

# ── Transition State Content ──────────────────────────────────────────────
ts_x = 5.6
transition_items = [
    ("Phase-In (Start):", "UAT rollout to all zones; Manager, CTO & Network AI modules enabled"),
    ("Phase-In (End):",   "Controlled production rollout complete; SLA tracking live"),
    ("Phase-Out (Start):","Legacy email/phone workflows sunset begins"),
    ("Phase-Out (End):",  "Spreadsheet trackers retired; full migration to TeleResolve"),
]
add_rect(slide4, ts_x, 1.7, phase_w, 4.8, fill=RGBColor(0xF0, 0xF5, 0xFF), line=KPMG_LBLUE, line_w=Pt(0.8))
for i, (title, desc) in enumerate(transition_items):
    ty = 1.8 + i * 1.15
    add_label(slide4, ts_x + 0.1, ty, phase_w - 0.2, 0.28, title,
              size=Pt(8.5), bold=True, color=KPMG_NAVY, align=PP_ALIGN.LEFT)
    add_label(slide4, ts_x + 0.1, ty + 0.27, phase_w - 0.2, 0.75, desc,
              size=Pt(7.5), color=KPMG_BLACK, align=PP_ALIGN.LEFT)

# ── Target State Content ──────────────────────────────────────────────────
tg_x = 8.8
target_items = [
    ("Phase-In (Start):", "Full deployment across all telecom circles; ML pipeline & Change Workflow fully live"),
    ("Phase-In (End):",   "Full business adoption; all 5 roles active across zones"),
    ("Phase-Out (Start):","Legacy manual processes fully retired"),
    ("Phase-Out (End):",  "100% digital, AI-powered complaint resolution; WhatsApp & Email alerts active"),
]
add_rect(slide4, tg_x, 1.7, phase_w, 4.8, fill=RGBColor(0xF0, 0xF5, 0xFF), line=KPMG_LBLUE, line_w=Pt(0.8))
for i, (title, desc) in enumerate(target_items):
    ty = 1.8 + i * 1.15
    add_label(slide4, tg_x + 0.1, ty, phase_w - 0.2, 0.28, title,
              size=Pt(8.5), bold=True, color=KPMG_NAVY, align=PP_ALIGN.LEFT)
    add_label(slide4, tg_x + 0.1, ty + 0.27, phase_w - 0.2, 0.75, desc,
              size=Pt(7.5), color=KPMG_BLACK, align=PP_ALIGN.LEFT)

# KPI targets summary bar at bottom
add_rect(slide4, 0.3, 6.5, 12.7, 0.45, fill=KPMG_NAVY, line=None)
kpi_targets = "Target KPIs:  ⬇ Complaint Resolution Time 60%  |  ⬆ CSAT Score > 85%  |  ⬆ SLA Compliance > 95%  |  ⬇ Manual Effort 80%  |  ⬆ First-Contact Resolution > 70%"
add_label(slide4, 0.5, 6.54, 12.3, 0.37, kpi_targets,
          size=Pt(8), bold=True, color=KPMG_CYAN, align=PP_ALIGN.CENTER)

# Footer
add_rect(slide4, 0, 7.1, 13.33, 0.4, fill=KPMG_NAVY)
add_label(slide4, 0.2, 7.15, 12, 0.25,
          "© 2026 TeleResolve. Confidential  |  Implementation Roadmap",
          size=Pt(8), color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)


# ── Save ────────────────────────────────────────────────────────────────────
output_path = r"c:\Users\didar\Downloads\files (12)\telecom-complaint-system\TeleResolve_Business_Architecture.pptx"
prs.save(output_path)
print(f"Saved → {output_path}")
