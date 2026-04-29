# -*- coding: utf-8 -*-
"""
TeleResolve - EXACT KPMG-format PPT Generator
Matches the exact visual style: swim lanes, X-circles, diamonds, ovals
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ── Palette ────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1F, 0x39, 0x64)
BLUE    = RGBColor(0x41, 0x72, 0xC4)
LBLUE   = RGBColor(0x9D, 0xB8, 0xE8)
CYAN    = RGBColor(0x00, 0xB0, 0xF0)
GREEN   = RGBColor(0x70, 0xAD, 0x47)
AMBER   = RGBColor(0xC5, 0x7C, 0x1E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x00, 0x00, 0x00)
LGREY   = RGBColor(0xF2, 0xF2, 0xF2)
DGREY   = RGBColor(0x40, 0x40, 0x40)
MGREY   = RGBColor(0xBF, 0xBF, 0xBF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Core draw helpers ──────────────────────────────────────────────────────
def shape(slide, mso_type, x, y, w, h, fill, line_rgb=MGREY, lw=Pt(0.75)):
    sp = slide.shapes.add_shape(mso_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line_rgb is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_rgb; sp.line.width = lw
    sp.text_frame.word_wrap = True
    return sp

def rect(slide, x, y, w, h, fill=BLUE, line=WHITE, lw=Pt(0.75)):
    return shape(slide, 1, x, y, w, h, fill, line, lw)   # 1=RECTANGLE

def rrect(slide, x, y, w, h, fill=BLUE, line=WHITE, lw=Pt(0.75)):
    sp = shape(slide, 5, x, y, w, h, fill, line, lw)     # 5=ROUNDED_RECT
    sp.adjustments[0] = 0.08
    return sp

def diamond(slide, x, y, w, h, fill=BLUE, line=WHITE):
    return shape(slide, 4, x, y, w, h, fill, line, Pt(0.75))  # 4=DIAMOND

def oval(slide, x, y, w, h, fill=NAVY, line=WHITE):
    return shape(slide, 9, x, y, w, h, fill, line, Pt(0.75))  # 9=OVAL

def xcircle(slide, cx, cy, r=0.18):
    """Small circle with X — junction connector, same as KPMG diagram"""
    sp = oval(slide, cx-r, cy-r, r*2, r*2, fill=DGREY, line=WHITE)
    txt(sp, "X", Pt(6), WHITE, bold=True)
    return sp

def txt(sp, text, size=Pt(7), color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    tf = sp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text; r.font.size = size
    r.font.bold = bold; r.font.color.rgb = color

def label(slide, x, y, w, h, text, size=Pt(7), color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold; r.font.color.rgb = color
    return tb

def arrow(slide, x1, y1, x2, y2, color=DGREY, w=Pt(1.0)):
    con = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    con.line.color.rgb = color; con.line.width = w
    ln = con._element.spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(con._element.spPr, qn('a:ln'))
    he = ln.find(qn('a:headEnd'))
    if he is None: he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type','arrow'); he.set('w','med'); he.set('len','med')
    te = ln.find(qn('a:tailEnd'))
    if te is None: te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type','none')

def lbl_arrow(slide, x, y, text, direction='h'):
    """Small Yes/No label near arrow"""
    label(slide, x, y, 0.25, 0.18, text, size=Pt(6.5), color=DGREY)

# ── Slide header / footer helpers ──────────────────────────────────────────
def header(slide, title, subtitle):
    rect(slide, 0, 0, 13.33, 0.52, fill=NAVY, line=None)
    label(slide, 0.18, 0.06, 5, 0.23, title,    size=Pt(13), color=WHITE, bold=True)
    label(slide, 0.18, 0.3,  8, 0.18, subtitle, size=Pt(8.5), color=CYAN)

def footer(slide, text):
    rect(slide, 0, 7.12, 13.33, 0.38, fill=NAVY, line=None)
    label(slide, 0.2, 7.15, 12, 0.22, text, size=Pt(7.5), color=MGREY)

def diagram_border(slide):
    rect(slide, 0.15, 0.6, 13.03, 6.55, fill=WHITE, line=MGREY, lw=Pt(1.2))

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
rect(s1, 0, 0, 13.33, 7.5, fill=NAVY, line=None)
rect(s1, 0, 0, 13.33, 0.07, fill=CYAN, line=None)
rect(s1, 0, 7.43, 13.33, 0.07, fill=CYAN, line=None)
label(s1, 0.4, 0.18, 2, 0.4, "KPMG", size=Pt(20), color=WHITE, bold=True)
label(s1, 1.2, 2.5, 10.9, 1.0,
      "TeleResolve\nTelecom Customer Complaint Handling System",
      size=Pt(28), color=WHITE, bold=True, align=PP_ALIGN.CENTER)
label(s1, 1.5, 3.7, 10, 0.45,
      "Business Architecture  |  As-Is Process  |  To-Be Process  |  Roadmap",
      size=Pt(13), color=CYAN, align=PP_ALIGN.CENTER)
rect(s1, 1.5, 4.3, 10.3, 0.03, fill=CYAN, line=None)
label(s1, 1.5, 4.45, 10, 0.28,
      "TeleResolve Team  |  April 2026  |  Confidential",
      size=Pt(10), color=LBLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AS-IS (matches KPMG As-Is exactly in layout)
# Single-flow diagram, no swim lanes, blue boxes on white
# ══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
rect(s2, 0, 0, 13.33, 7.5, fill=WHITE, line=None)
header(s2, "6. Business Architecture", "Business Process (As-Is) — Manual Complaint Handling (TeleResolve)")
diagram_border(s2)

# ── TOP ROW ────────────────────────────────────────────────────────────────
# Start → Complaint Info Available? → Yes → Collect Customer Data → Run Manual Triage
bx, by, bw, bh = 0.3, 0.14, 0.95, 0.42   # box width/height offsets

# Start oval
sp = oval(s2, 0.35, 1.05, 0.68, 0.42, fill=NAVY, line=WHITE)
txt(sp, "Start", Pt(7.5), WHITE, bold=True)
arrow(s2, 1.03, 1.26, 1.45, 1.26)

# Decision: Complaint Details Available?
sp = diamond(s2, 1.45, 0.92, 1.1, 0.68, fill=BLUE, line=WHITE)
txt(sp, "Complaint\nDetails\nAvailable?", Pt(6.5), WHITE)
lbl_arrow(s2, 2.58, 1.12, "Yes")
arrow(s2, 2.55, 1.26, 2.98, 1.26)

# Yes path: Collect Customer Data
sp = rrect(s2, 2.98, 1.04, 1.15, 0.42, fill=BLUE, line=WHITE)
txt(sp, "Collect\nCustomer Data", Pt(6.5), WHITE)
arrow(s2, 4.13, 1.25, 4.55, 1.25)

# Run Manual Triage
sp = rrect(s2, 4.55, 1.04, 1.15, 0.42, fill=BLUE, line=WHITE)
txt(sp, "Run Manual\nTriage", Pt(6.5), WHITE)
arrow(s2, 5.7, 1.25, 9.8, 1.25)

# No path from Decision (goes down)
lbl_arrow(s2, 1.35, 1.63, "No")
arrow(s2, 1.99, 1.6, 1.99, 2.05)

# ── MIDDLE FLOW (No path) ─────────────────────────────────────────────────
# Request Customer Info/Evidence
sp = rrect(s2, 1.45, 2.05, 1.15, 0.42, fill=BLUE, line=WHITE)
txt(sp, "Request\nCustomer Info\n/ Evidence", Pt(6.0), WHITE)
arrow(s2, 2.6, 2.26, 3.02, 2.26)

# Receive Complaint Info
sp = rrect(s2, 3.02, 2.05, 1.15, 0.42, fill=BLUE, line=WHITE)
txt(sp, "Receive\nComplaint Info\n/ Call Log", Pt(6.0), WHITE)
arrow(s2, 4.17, 2.26, 4.59, 2.26)

# Identify Issue
sp = rrect(s2, 4.59, 2.05, 1.15, 0.42, fill=BLUE, line=WHITE)
txt(sp, "Identify\nIssue\nManually", Pt(6.0), WHITE)
arrow(s2, 5.74, 2.26, 6.16, 2.26)

# Decision: Escalation Required?
sp = diamond(s2, 6.16, 1.92, 1.1, 0.68, fill=BLUE, line=WHITE)
txt(sp, "Escalation\n/ Follow-up\nRequired?", Pt(6.0), WHITE)
lbl_arrow(s2, 7.29, 2.1, "Yes")
arrow(s2, 7.26, 2.26, 7.68, 2.26)

# Escalation / Remediation
sp = rrect(s2, 7.68, 2.05, 1.15, 0.42, fill=BLUE, line=WHITE)
txt(sp, "Escalation &\nRemediation\nTracking", Pt(6.0), WHITE)
arrow(s2, 8.83, 2.26, 9.25, 2.26)

# Reviewer / Manager Review (top right)
sp = rrect(s2, 9.25, 1.04, 1.15, 0.42, fill=BLUE, line=WHITE)
txt(sp, "Reviewer /\nManager\nReview", Pt(6.0), WHITE)
# connect escalation up to manager review
arrow(s2, 9.83, 2.05, 9.83, 1.46)

# Right side: Finalize & Close Ticket
sp = rrect(s2, 10.6, 1.04, 1.15, 0.42, fill=BLUE, line=WHITE)
txt(sp, "Finalize &\nClose Ticket\nManually", Pt(6.0), WHITE)
arrow(s2, 10.4, 1.25, 10.6, 1.25)
arrow(s2, 11.75, 1.25, 12.12, 1.25)

# End oval
sp = oval(s2, 12.12, 1.05, 0.68, 0.42, fill=NAVY, line=WHITE)
txt(sp, "Stop", Pt(7.5), WHITE, bold=True)

# ── BOTTOM BOX: Manual Process note ──────────────────────────────────────
sp = rrect(s2, 3.02, 3.1, 3.5, 0.55, fill=BLUE, line=WHITE)
txt(sp, "Manual Categorization / Priority Assignment / SLA Tracking / Resolution Steps", Pt(6.5), WHITE)

# connect Identify Issue down to this box
arrow(s2, 5.16, 2.47, 5.16, 3.1)
arrow(s2, 5.16, 3.65, 5.16, 2.47)  # back up

# ── PAIN POINTS box (right side, matches KPMG style) ─────────────────────
sp = rrect(s2, 0.25, 3.75, 12.83, 2.95, fill=RGBColor(0xDA,0xE3,0xF3), line=MGREY)
label(s2, 0.4,  3.82, 3, 0.28, "Key Pain Points (As-Is)", size=Pt(9), color=NAVY, bold=True)
pains = [
    ("No unified complaint portal — customers call/email manually",
     "No AI-assisted diagnosis — engineers investigate manually"),
    ("Ticket creation via spreadsheets — no real-time tracking",
     "No SLA alerts — breaches go unnoticed"),
    ("Manual escalation via email/phone — slow response",
     "No change request audit trail — RF changes undocumented"),
    ("No network KPI visibility for agents",
     "No CSAT collection — no feedback loop"),
    ("No geospatial view of impacted sites",
     "No analytics dashboard for CTO decision-making"),
]
for i, (left, right) in enumerate(pains):
    ry = 4.15 + i * 0.47
    label(s2, 0.4,  ry, 6.2, 0.4, "• " + left,  size=Pt(7.5), color=DGREY)
    label(s2, 6.75, ry, 6.2, 0.4, "• " + right, size=Pt(7.5), color=DGREY)

footer(s2, "© 2026 KPMG Assurance and Consulting Services LLP  |  Business Architecture — As-Is Process  |  Confidential")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TO-BE  (matches KPMG To-Be exactly — 3 swim lanes + X circles)
# ══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
rect(s3, 0, 0, 13.33, 7.5, fill=WHITE, line=None)
header(s3, "6. Business Architecture", "Business Process (To-Be) — TeleResolve AI-Powered System")
diagram_border(s3)

LANE_LBL_W = 1.1
LANE_X = 0.18

# ── Draw swim lane bands ───────────────────────────────────────────────────
lane_defs = [
    ("Customer",          0.65, 1.5),
    ("Admin\nApproval",   2.2,  1.1),
    ("TeleResolve\nSystem",3.35, 3.35),
]
for lname, ly, lh in lane_defs:
    # label box
    sp = rect(s3, LANE_X, ly, LANE_LBL_W, lh, fill=NAVY, line=WHITE, lw=Pt(0.5))
    txt(sp, lname, Pt(8), WHITE, bold=True)
    # lane BG
    bg = RGBColor(0xEE,0xF3,0xFB) if lane_defs.index((lname,ly,lh))%2==0 else WHITE
    rect(s3, LANE_X+LANE_LBL_W, ly, 13.33-LANE_X-LANE_LBL_W-0.15, lh, fill=bg, line=MGREY, lw=Pt(0.4))

# ── LANE 1: CUSTOMER ───────────────────────────────────────────────────────
L1Y = 0.82   # vertical center-ish of lane 1

# Start oval
sp = oval(s3, 1.35, L1Y, 0.7, 0.42, fill=NAVY)
txt(sp, "Start", Pt(7.5), WHITE, bold=True)
arrow(s3, 2.05, L1Y+0.21, 2.45, L1Y+0.21)

# Decision: Existing User?
sp = diamond(s3, 2.45, L1Y-0.1, 1.0, 0.62, fill=BLUE)
txt(sp, "Existing\nUser?", Pt(6.5), WHITE)

# No → Register
lbl_arrow(s3, 3.48, L1Y+0.08, "No")
arrow(s3, 3.45, L1Y+0.21, 3.85, L1Y+0.21)
sp = rrect(s3, 3.85, L1Y, 1.0, 0.42, fill=BLUE)
txt(sp, "Register\n(Web / App)", Pt(6.5), WHITE)
arrow(s3, 4.85, L1Y+0.21, 5.25, L1Y+0.21)

# Create Password
sp = rrect(s3, 5.25, L1Y, 1.0, 0.42, fill=BLUE)
txt(sp, "Create\nPassword", Pt(6.5), WHITE)
arrow(s3, 6.25, L1Y+0.21, 6.65, L1Y+0.21)

# Generate OTP
sp = rrect(s3, 6.65, L1Y, 1.0, 0.42, fill=BLUE)
txt(sp, "Generate\nOTP", Pt(6.5), WHITE)
arrow(s3, 7.65, L1Y+0.21, 10.3, L1Y+0.21)

# Yes path from diamond directly to Login (arc over top)
lbl_arrow(s3, 2.85, L1Y-0.14, "Yes")
# draw top arc line: diamond top → straight line right → Login
arrow(s3, 2.95, L1Y-0.1, 2.95, L1Y-0.35)
arrow(s3, 2.95, L1Y-0.35, 10.65, L1Y-0.35)
arrow(s3, 10.65, L1Y-0.35, 10.65, L1Y)

# Login box (far right of lane 1)
sp = rrect(s3, 10.3, L1Y, 1.0, 0.42, fill=NAVY)
txt(sp, "Login", Pt(8), WHITE, bold=True)

# ── LANE 2: ADMIN APPROVAL ─────────────────────────────────────────────────
L2Y = 2.42

# Arrow from Register down to Admin lane
arrow(s3, 4.35, L1Y+0.42, 4.35, L2Y+0.04)

# Admin Approval diamond
sp = diamond(s3, 3.85, L2Y, 1.0, 0.6, fill=BLUE)
txt(sp, "Admin\nApproval?", Pt(6.5), WHITE)
# No → back to Register area
lbl_arrow(s3, 3.3, L2Y+0.28, "No")
arrow(s3, 3.37, L2Y+0.3, 1.5, L2Y+0.3)
sp = rrect(s3, 1.3, L2Y+0.1, 1.0, 0.42, fill=RGBColor(0xC0,0x00,0x00))
txt(sp, "Reject &\nNotify", Pt(6.5), WHITE)
# Yes → continue to OTP
lbl_arrow(s3, 4.88, L2Y+0.14, "Yes")
arrow(s3, 4.85, L2Y+0.3, 6.15, L2Y+0.3)
arrow(s3, 6.15, L2Y+0.3, 6.15, L1Y+0.42)   # back up to OTP flow

# ── LANE 3: TELERESOLVE SYSTEM (3 rows, matches Control Iris Tool layout) ──
L3_TOP  = 3.45    # top of lane 3
L3_R1Y  = 3.57    # row 1 y
L3_R2Y  = 4.42    # row 2 y
L3_R3Y  = 5.27    # row 3 y
BOX_H   = 0.42
BOX_W   = 1.05

# Arrow from Login down into TeleResolve System lane
arrow(s3, 10.8, L1Y+0.42, 10.8, L3_R1Y)

# ── ROW 1 (top row of TeleResolve lane) ────────────────────────────────────
# X junction left
xcircle(s3, 1.62, L3_R1Y+BOX_H/2)
arrow(s3, 1.8, L3_R1Y+BOX_H/2, 2.2, L3_R1Y+BOX_H/2)

r1_items = [
    "Submit\nComplaint",
    "View Active\nTickets",
    "Recent\nSessions",
    "Network AI\nChat",
    "Upload\nEvidence",
]
cx = 2.2
for i, it in enumerate(r1_items):
    sp = rrect(s3, cx, L3_R1Y, BOX_W, BOX_H, fill=BLUE)
    txt(sp, it, Pt(6.5), WHITE)
    if i < len(r1_items)-1:
        arrow(s3, cx+BOX_W, L3_R1Y+BOX_H/2, cx+BOX_W+0.1, L3_R1Y+BOX_H/2)
        cx += BOX_W + 0.1
    else:
        cx += BOX_W

# X junction right
xcircle(s3, cx+0.18, L3_R1Y+BOX_H/2)

# ── ROW 2 (middle row) ─────────────────────────────────────────────────────
sp = rrect(s3, 4.5, L3_R2Y, 1.35, BOX_H, fill=BLUE)
txt(sp, "Select Issue\nCategory", Pt(6.5), WHITE)
arrow(s3, 5.85, L3_R2Y+BOX_H/2, 6.25, L3_R2Y+BOX_H/2)
sp = rrect(s3, 6.25, L3_R2Y, 1.35, BOX_H, fill=BLUE)
txt(sp, "View\nDashboard", Pt(6.5), WHITE)

# Arrows connecting X-right (row1) down → Select Issue Category (row2)
arrow(s3, cx+0.18, L3_R1Y+BOX_H, cx+0.18, L3_R2Y+BOX_H/2)
arrow(s3, cx+0.18, L3_R2Y+BOX_H/2, 4.5, L3_R2Y+BOX_H/2)

# ── ROW 3 (bottom row) ─────────────────────────────────────────────────────
# X junction left row3
xcircle(s3, 1.62, L3_R3Y+BOX_H/2)
# connect left X row1 down to X row3
arrow(s3, 1.62, L3_R1Y+BOX_H, 1.62, L3_R3Y)

arrow(s3, 1.8, L3_R3Y+BOX_H/2, 2.2, L3_R3Y+BOX_H/2)

r3_items = [
    ("Start New\nChat Session",  BLUE),
    ("AI Triage &\nResolution",  BLUE),
    ("Run Network\nDiagnosis",   BLUE),
    ("Live\nDashboard",          BLUE),
]
cx3 = 2.2
for i, (it, fc) in enumerate(r3_items):
    sp = rrect(s3, cx3, L3_R3Y, BOX_W, BOX_H, fill=fc)
    txt(sp, it, Pt(6.5), WHITE)
    if i < len(r3_items)-1:
        # X circle between boxes (like KPMG diagram)
        nx = cx3 + BOX_W + 0.18
        xcircle(s3, cx3+BOX_W+0.18, L3_R3Y+BOX_H/2)
        arrow(s3, cx3+BOX_W, L3_R3Y+BOX_H/2, cx3+BOX_W+0.08, L3_R3Y+BOX_H/2)
        arrow(s3, cx3+BOX_W+0.28, L3_R3Y+BOX_H/2, cx3+BOX_W+0.36, L3_R3Y+BOX_H/2)
        cx3 += BOX_W + 0.36
    else:
        cx3 += BOX_W

arrow(s3, cx3, L3_R3Y+BOX_H/2, cx3+0.15, L3_R3Y+BOX_H/2)
# End oval
sp = oval(s3, cx3+0.15, L3_R3Y, 0.7, BOX_H, fill=NAVY)
txt(sp, "End", Pt(7.5), WHITE, bold=True)

# Login → row1 left X connector
arrow(s3, 10.8, L3_R1Y, 10.8, L3_R1Y)
# connect view dashboard to Live Dashboard row3
arrow(s3, 7.6, L3_R2Y+BOX_H/2, 10.5, L3_R2Y+BOX_H/2)
arrow(s3, 10.5, L3_R2Y+BOX_H/2, 10.5, L3_R3Y+BOX_H/2)

footer(s3, "© 2026 KPMG Assurance and Consulting Services LLP  |  Business Architecture — To-Be Process  |  Confidential")


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ROADMAP  (exact KPMG layout: amber platform box + 3 chevrons)
# ══════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
rect(s4, 0, 0, 13.33, 7.5, fill=WHITE, line=None)
header(s4, "7. Roadmap", "TeleResolve Platform Implementation Roadmap")
diagram_border(s4)

# Outer roadmap box
rect(s4, 0.3, 0.7, 12.73, 6.55, fill=LGREY, line=MGREY, lw=Pt(1))

# Amber platform box (top-left, exactly like KPMG)
sp = rrect(s4, 0.5, 1.0, 1.55, 0.92, fill=AMBER, line=WHITE)
txt(sp, "TeleResolve\nPlatform", Pt(10), WHITE, bold=True)

# Three chevron arrows (green, side by side)
CHEV_Y  = 0.95
CHEV_H  = 0.52
CHEV_W  = 3.1
GAP     = 0.15
CHEV_X  = [2.3, 2.3+CHEV_W+GAP, 2.3+2*(CHEV_W+GAP)]

phase_titles = ["Current State", "Transition State", "Target State"]
for i, (px, pt) in enumerate(zip(CHEV_X, phase_titles)):
    sp = rrect(s4, px, CHEV_Y, CHEV_W, CHEV_H, fill=GREEN, line=WHITE, lw=Pt(1.5))
    sp.adjustments[0] = 0.03
    txt(sp, pt, Pt(11), WHITE, bold=True)
    if i < 2:
        arrow(s4, px+CHEV_W, CHEV_Y+CHEV_H/2, px+CHEV_W+GAP, CHEV_Y+CHEV_H/2,
              color=GREEN, w=Pt(2))

# Content under each phase
CONT_Y  = 1.58   # top of content area
CONT_H  = 5.5
phase_content = [
    [
        ("Phase-In (Start):",  "Pilot launch for 1 telecom zone; Customer chatbot, ticket management & agent dashboard go live"),
        ("Phase-In (End):",    "Pilot readiness confirmed; Customer & Agent roles fully tested end-to-end"),
        ("Phase-Out (Start):", "Manual spreadsheet-based complaint logging begins phased reduction"),
        ("Phase-Out (End):",   "Pilot zone fully migrated to TeleResolve; legacy spreadsheets retired"),
    ],
    [
        ("Phase-In (Start):",  "UAT rollout to all zones; Manager approvals, CTO dashboards & Network AI modules enabled"),
        ("Phase-In (End):",    "Controlled production rollout complete; SLA tracking & WhatsApp alerts live"),
        ("Phase-Out (Start):", "Legacy phone/email escalation channels sunset begins"),
        ("Phase-Out (End):",   "All spreadsheet trackers retired; full migration to TeleResolve complete"),
    ],
    [
        ("Phase-In (Start):",  "Full deployment across all telecom circles; ML pipeline, Change Workflow & KPI analytics live"),
        ("Phase-In (End):",    "Full business adoption; all 5 roles active — Customer, Agent, Manager, CTO, Admin"),
        ("Phase-Out (Start):", "All legacy manual processes retired system-wide"),
        ("Phase-Out (End):",   "100% digital AI-powered complaint resolution; all notifications automated"),
    ],
]

for i, (px, items) in enumerate(zip(CHEV_X, phase_content)):
    # content border box
    rect(s4, px, CONT_Y, CHEV_W, CONT_H, fill=WHITE, line=MGREY, lw=Pt(0.6))
    for j, (bold_txt, desc_txt) in enumerate(items):
        iy = CONT_Y + 0.12 + j * 1.28
        # bullet dot
        rect(s4, px+0.12, iy+0.07, 0.06, 0.06, fill=NAVY, line=None)
        label(s4, px+0.22, iy, CHEV_W-0.28, 0.25, bold_txt,
              size=Pt(8), color=NAVY, bold=True)
        label(s4, px+0.22, iy+0.24, CHEV_W-0.28, 0.92, desc_txt,
              size=Pt(7.5), color=DGREY)

# KPI targets bar at bottom
rect(s4, 0.3, 6.88, 12.73, 0.37, fill=NAVY, line=None)
label(s4, 0.5, 6.91, 12.3, 0.3,
      "Target Outcomes:   Resolution Time -60%   |   CSAT > 85%   |   SLA Compliance > 95%   |   Manual Effort -80%   |   FCR > 70%",
      size=Pt(8), color=CYAN, bold=True, align=PP_ALIGN.CENTER)

footer(s4, "© 2026 KPMG Assurance and Consulting Services LLP  |  Roadmap  |  Confidential")


# ── Save ───────────────────────────────────────────────────────────────────
out = r"c:\Users\didar\Downloads\files (12)\telecom-complaint-system\TeleResolve_Business_Architecture_v2.pptx"
prs.save(out)
print("Saved: " + out)
