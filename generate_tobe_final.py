# -*- coding: utf-8 -*-
"""
TeleResolve - To-Be Diagram
EXACT replica of KPMG format:
  Lane 1 (Customer): Start -> Existing User? --No--> Register -> Create Password -> Generate OTP -> Login
                                              --Yes (top arc)------------------------------------> Login
  Lane 2 (Admin Approval): diamond --Yes/No-->
  Lane 3 (TeleResolve System):
    Row 1: [Left-X] -> 5 boxes -> [Right-X]
    Row 2:             Select Issue Category -> View Dashboard
    Row 3: [Start AI Chat oval] -> [X] -> Run Network Diagnosis -> [X] -> Live Dashboard -> End
  All X circles identical to original. Only box text changed for TeleResolve.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ── Colours (matching KPMG blue palette) ──────────────────────────────────
NAVY  = RGBColor(0x1F, 0x39, 0x64)
BLUE  = RGBColor(0x41, 0x72, 0xC4)
CYAN  = RGBColor(0x00, 0xB0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY  = RGBColor(0xBF, 0xBF, 0xBF)
DGREY = RGBColor(0x35, 0x35, 0x35)
LBBL  = RGBColor(0xEC, 0xF2, 0xFB)
WBBL  = RGBColor(0xF8, 0xFB, 0xFF)
RED   = RGBColor(0xC0, 0x00, 0x00)

# ── Presentation ───────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Primitive helpers ──────────────────────────────────────────────────────
def _sp(mso, x, y, w, h, fill, lc=GREY, lw=Pt(0.75)):
    s = sl.shapes.add_shape(mso, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: s.fill.background()
    else: s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc is None: s.line.fill.background()
    else: s.line.color.rgb = lc; s.line.width = lw
    s.text_frame.word_wrap = True
    return s

def RCT(x,y,w,h, fill=BLUE, lc=WHITE, lw=Pt(0.75)):  return _sp(1,x,y,w,h,fill,lc,lw)
def RND(x,y,w,h, fill=BLUE, lc=WHITE, lw=Pt(0.75)):   # rounded rect
    s=_sp(5,x,y,w,h,fill,lc,lw); s.adjustments[0]=0.09; return s
def DMD(x,y,w,h, fill=BLUE, lc=WHITE):                return _sp(4,x,y,w,h,fill,lc,Pt(0.75))
def OVL(x,y,w,h, fill=NAVY, lc=WHITE):                return _sp(9,x,y,w,h,fill,lc,Pt(0.75))

def XC(cx, cy, r=0.175):
    """Dark circle with X — junction connector identical to KPMG original."""
    s = OVL(cx-r, cy-r, r*2, r*2, fill=DGREY, lc=WHITE)
    T(s, "X", Pt(6.5), WHITE, bold=True); return s

def T(sp, text, sz=Pt(7), col=WHITE, bold=False, align=PP_ALIGN.CENTER):
    tf=sp.text_frame; p=tf.paragraphs[0]; p.alignment=align
    r=p.runs[0] if p.runs else p.add_run()
    r.text=text; r.font.size=sz; r.font.bold=bold; r.font.color.rgb=col

def LBL(x,y,w,h, text, sz=Pt(7), col=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.size=sz; r.font.bold=bold; r.font.color.rgb=col

def ARR(x1,y1,x2,y2, col=DGREY, w=Pt(1.0)):
    """Straight arrow connector from (x1,y1) to (x2,y2) in inches."""
    c=sl.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=col; c.line.width=w
    ln=c._element.spPr.find(qn('a:ln'))
    if ln is None: ln=etree.SubElement(c._element.spPr,qn('a:ln'))
    he=ln.find(qn('a:headEnd'))
    if he is None: he=etree.SubElement(ln,qn('a:headEnd'))
    he.set('type','arrow'); he.set('w','med'); he.set('len','med')
    te=ln.find(qn('a:tailEnd'))
    if te is None: te=etree.SubElement(ln,qn('a:tailEnd'))
    te.set('type','none')

# ═══════════════════════════════════════════════════════════════════════════
# LAYOUT  (all in inches)
# ═══════════════════════════════════════════════════════════════════════════
LLW = 1.1     # lane-label column width
LX  = 0.15   # diagram left  x
RX  = 13.18  # diagram right x
CX  = LX+LLW  # content area starts here = 1.25
BH  = 0.44   # standard box height
XR  = 0.175  # X-circle radius

# Swim lane y-ranges (exactly proportional to KPMG image)
L1Y, L1H = 0.62, 1.50   # Customer         → 0.62 – 2.12
L2Y, L2H = 2.12, 1.08   # Admin Approval   → 2.12 – 3.20
L3Y, L3H = 3.20, 3.82   # TeleResolve Sys  → 3.20 – 7.02

# Vertical midpoints for each row in lane 3
R1M = L3Y + 0.70         # row-1 box mid-y  3.90
R2M = L3Y + 1.95         # row-2 box mid-y  5.15
R3M = L3Y + 3.20         # row-3 box mid-y  6.40

R1T = R1M - BH/2         # box top-y row 1  3.68
R2T = R2M - BH/2         # box top-y row 2  4.93
R3T = R3M - BH/2         # box top-y row 3  6.18

# Lane-1 box mid-y
L1M = L1Y + L1H/2        # 1.37
L1T = L1M - BH/2         # 1.15

# Lane-2 diamond center-y
L2M = L2Y + L2H/2        # 2.66

# ── Row-1: 5 boxes between Left-X and Right-X ─────────────────────────────
XL1X = CX + 0.20         # left-X  center-x  1.45
XR1X = RX - 0.25         # right-X center-x 12.93

BOX_AREA = XR1X - XR - 0.12 - (XL1X + XR + 0.12)   # usable x for boxes
# = 12.93-0.175-0.12 - (1.45+0.175+0.12) = 12.635 - 1.745 = 10.89
GAP  = 0.14
BW1  = (BOX_AREA - 4*GAP) / 5                        # ≈ 2.09
B1X  = XL1X + XR + 0.12                              # first box left  1.745

# Row-1 box left-x positions
R1XS = [B1X + i*(BW1+GAP) for i in range(5)]

# Row-2 boxes (centered under the diagram, left-to-right flow)
SIC_W, VD_W = 1.65, 1.65
SIC_X = 4.70                  # Select Issue Category
VD_X  = SIC_X + SIC_W + 0.18  # View Dashboard

# Row-3 items
# Start AI Chat oval (left, at row-3 level, enters from left-X vertical drop)
SAC_W = 1.0
SAC_X = XL1X + XR + 0.12     # just right of left-X  ≈ 1.745
# X1 between Start-AI-Chat and Run Diagnosis
X1_CX = SAC_X + SAC_W + 0.20 + XR
# Run Network Diagnosis box
RND_X = X1_CX + XR + 0.12
RND_W = 1.55
# X2 – merge junction (View Dashboard drops into here)
X2_CX = VD_X + VD_W/2        # aligned with View Dashboard center  ≈ 7.60
# Live Dashboard
LD_X  = X2_CX + XR + 0.18
LD_W  = 1.50
# End oval
END_X = LD_X + LD_W + 0.18
END_W = 0.75

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE CHROME
# ═══════════════════════════════════════════════════════════════════════════
RCT(0,0,13.33,7.5, fill=WHITE, lc=None)  # white bg

# Header (navy bar)
RCT(0, 0, 13.33, 0.52, fill=NAVY, lc=None)
LBL(0.18, 0.06, 6, 0.22, "6. Business Architecture",
    sz=Pt(13), col=WHITE, bold=True)
LBL(0.18, 0.30, 10, 0.18,
    "Business Process (To-Be)  —  TeleResolve AI-Powered Complaint Handling System",
    sz=Pt(8.5), col=CYAN)

# Footer
RCT(0, 7.12, 13.33, 0.38, fill=NAVY, lc=None)
LBL(0.2, 7.14, 12.9, 0.22,
    "© 2026 KPMG Assurance and Consulting Services LLP an Indian Limited Liability Partnership and a member firm of the KPMG global organization of independent member firms affiliated with KPMG International Limited, a private English company limited by guarantee. All rights reserved.",
    sz=Pt(6), col=GREY)

# Diagram outer border
RCT(LX, 0.60, RX-LX, L3Y+L3H-0.60, fill=WHITE, lc=GREY, lw=Pt(1.0))

# ═══════════════════════════════════════════════════════════════════════════
# SWIM LANE BANDS
# ═══════════════════════════════════════════════════════════════════════════
for name, ly, lh, bg in [
    ("Customer",            L1Y, L1H, LBBL),
    ("Admin\nApproval",     L2Y, L2H, WBBL),
    ("TeleResolve\nSystem", L3Y, L3H, LBBL),
]:
    sp = RCT(LX, ly, LLW, lh, fill=NAVY, lc=WHITE, lw=Pt(0.5))
    T(sp, name, sz=Pt(8.5), col=WHITE, bold=True)
    RCT(CX, ly, RX-CX, lh, fill=bg, lc=GREY, lw=Pt(0.5))

# ═══════════════════════════════════════════════════════════════════════════
# LANE 1 — CUSTOMER
# ═══════════════════════════════════════════════════════════════════════════
# Start oval
sp = OVL(CX+0.05, L1T, 0.72, BH); T(sp, "Start", Pt(8), WHITE, bold=True)
S_CX = CX + 0.05 + 0.36   # start oval center-x  1.61
ARR(S_CX+0.36, L1M, S_CX+0.36+0.18, L1M)   # → diamond

# Diamond: Existing User?
DW, DH = 1.05, 0.68
DX = S_CX + 0.36 + 0.18   # diamond left-x  2.15
D_CX = DX + DW/2           # diamond center-x  2.675
D_CY = L1M                 # diamond center-y  1.37
D_TOP= L1T - (DH-BH)/2     # diamond top-y   1.02
sp = DMD(DX, D_TOP, DW, DH); T(sp, "Existing\nUser?", Pt(6.5), WHITE)

# YES arc — above lane 1 boxes, goes over the top to Login
ARC_Y = L1Y + 0.07          # arc height near top of lane 1   0.69
LBL(D_CX + 0.04, ARC_Y-0.02, 0.3, 0.18, "Yes", sz=Pt(6.5), col=DGREY)
ARR(D_CX, D_TOP,  D_CX, ARC_Y)        # up from diamond top
ARR(D_CX, ARC_Y,  12.30, ARC_Y)       # right across the top
ARR(12.30, ARC_Y, 12.30, L1T)         # down to Login top

# NO → right → Register
NO_X = DX + DW              # right edge of diamond  3.20
LBL(NO_X+0.04, L1M-0.14, 0.25, 0.18, "No", sz=Pt(6.5), col=DGREY)
ARR(NO_X, L1M, NO_X+0.25, L1M)

# Register
REG_X = NO_X + 0.25
REG_W = 1.05
sp = RND(REG_X, L1T, REG_W, BH); T(sp, "Register", Pt(7), WHITE)
REG_CX = REG_X + REG_W/2    # register center-x  3.875
REG_BOT = L1T + BH           # register bottom-y  1.59
ARR(REG_X+REG_W, L1M, REG_X+REG_W+0.18, L1M)

# Create Password
CP_X = REG_X + REG_W + 0.18
CP_W = 1.20
sp = RND(CP_X, L1T, CP_W, BH); T(sp, "Create\nPassword", Pt(6.5), WHITE)
ARR(CP_X+CP_W, L1M, CP_X+CP_W+0.18, L1M)

# Generate OTP
OTP_X = CP_X + CP_W + 0.18
OTP_W = 1.20
sp = RND(OTP_X, L1T, OTP_W, BH); T(sp, "Generate\nOTP", Pt(6.5), WHITE)
OTP_R = OTP_X + OTP_W       # right edge of OTP  7.41
ARR(OTP_R, L1M, 11.75, L1M)  # long arrow → Login

# Login (far right)
LOGIN_X = 11.75
LOGIN_W = 1.0
LOGIN_CX = LOGIN_X + LOGIN_W/2  # 12.25
sp = RND(LOGIN_X, L1T, LOGIN_W, BH, fill=NAVY); T(sp, "Login", Pt(8.5), WHITE, bold=True)
LOGIN_BOT = L1T + BH            # bottom of Login  1.59

# ═══════════════════════════════════════════════════════════════════════════
# LANE 2 — ADMIN APPROVAL
# ═══════════════════════════════════════════════════════════════════════════
ADW, ADH = 1.05, 0.70
AD_CX = REG_CX               # align under Register center  3.875
ADX = AD_CX - ADW/2          # diamond left-x  3.35
AD_TOP = L2M - ADH/2         # diamond top  2.31
sp = DMD(ADX, AD_TOP, ADW, ADH); T(sp, "Admin\nApproval", Pt(6.5), WHITE)

# Arrow: Register bottom → Admin diamond top
ARR(REG_CX, REG_BOT, REG_CX, AD_TOP)

# NO → left → Reject box
LBL(ADX - 0.28, L2M - 0.10, 0.25, 0.18, "No", sz=Pt(6.5), col=DGREY)
ARR(ADX, L2M, CX+0.10, L2M)
sp = RND(CX+0.10, L2M-BH/2, 0.95, BH, fill=RED); T(sp, "Reject &\nNotify", Pt(6.5), WHITE)

# YES → right → up to Login bottom
LBL(ADX+ADW+0.04, L2M-0.10, 0.28, 0.18, "Yes", sz=Pt(6.5), col=DGREY)
ARR(ADX+ADW, L2M, LOGIN_CX, L2M)   # right to Login center-x
ARR(LOGIN_CX, L2M, LOGIN_CX, LOGIN_BOT)  # up to Login bottom

# ═══════════════════════════════════════════════════════════════════════════
# LANE 3 — TELERESOLVE SYSTEM
# ═══════════════════════════════════════════════════════════════════════════

# Login → drop into lane 3 top → travel left → Left-X circle
ARR(LOGIN_CX, LOGIN_BOT, LOGIN_CX, L3Y+0.22)   # down from Login
ARR(LOGIN_CX, L3Y+0.22, XL1X, L3Y+0.22)         # left along top of lane 3
ARR(XL1X, L3Y+0.22, XL1X, R1M-XR)              # down into Left-X

# ── Left-X circle (row 1 entry & row 3 vertical feed) ─────────────────────
XC(XL1X, R1M)

# Left-X → row 1 first box
ARR(XL1X+XR, R1M, R1XS[0], R1M)

# ── ROW 1: 5 feature boxes ────────────────────────────────────────────────
ROW1_LABELS = [
    "Submit\nComplaint",
    "View Active\nTickets",
    "Recent\nSessions",
    "Network AI\nChat",
    "Upload\nEvidence",
]
for i, (rx0, lbl_txt) in enumerate(zip(R1XS, ROW1_LABELS)):
    sp = RND(rx0, R1T, BW1, BH); T(sp, lbl_txt, Pt(6.5), WHITE)
    if i < 4:
        ARR(rx0+BW1, R1M, R1XS[i+1], R1M)

# last box → Right-X
ARR(R1XS[4]+BW1, R1M, XR1X-XR, R1M)

# Right-X circle
XC(XR1X, R1M)

# ── ROW 2: Select Issue Category → View Dashboard ─────────────────────────
# Right-X drops down to row-2, travels left to reach Select Issue Category
ARR(XR1X, R1M+XR, XR1X, R2M)           # down from Right-X
ARR(XR1X, R2M, SIC_X+SIC_W, R2M)        # left to right-edge of View Dashboard
                                          # (enters from the right: Right-X → VD → SIC)
# View Dashboard
sp = RND(VD_X, R2T, VD_W, BH); T(sp, "View\nDashboard", Pt(6.5), WHITE)
ARR(VD_X, R2M, SIC_X+SIC_W, R2M)        # VD left ← right from Right-X path (above already done)
# Actually flow is right-to-left from Right-X down to VD then to SIC:
# right-X side: already drew down+left to SIC right edge
# Arrow SIC→VD (left to right inside row 2):
sp = RND(SIC_X, R2T, SIC_W, BH); T(sp, "Select Issue\nCategory", Pt(6.5), WHITE)
ARR(SIC_X+SIC_W, R2M, VD_X, R2M)        # Select → View Dashboard

# View Dashboard drops down to X2 merge
VD_CX = VD_X + VD_W/2
ARR(VD_CX, R2T+BH, VD_CX, R3M-XR)      # straight down to X2

# ── Left-X vertical drop to Row 3 ─────────────────────────────────────────
ARR(XL1X, R1M+XR, XL1X, R3M)            # down from Left-X to row-3 level
ARR(XL1X, R3M, SAC_X, R3M)              # short right to Start AI Chat left edge

# ── ROW 3: Start AI Chat → X1 → Run Network Diagnosis → X2 → Live Dashboard → End ──
# Start AI Chat oval
sp = OVL(SAC_X, R3T, SAC_W, BH, fill=NAVY); T(sp, "Start AI\nChat", Pt(6.5), WHITE)
ARR(SAC_X+SAC_W, R3M, X1_CX-XR, R3M)   # → X1

# X1 junction circle
XC(X1_CX, R3M)
ARR(X1_CX+XR, R3M, RND_X, R3M)          # X1 → Run Network Diagnosis

# Run Network Diagnosis
sp = RND(RND_X, R3T, RND_W, BH); T(sp, "Run Network\nDiagnosis", Pt(6.5), WHITE)
ARR(RND_X+RND_W, R3M, X2_CX-XR, R3M)   # → X2

# X2 junction circle (merge: Row-3 flow + View Dashboard from above)
XC(X2_CX, R3M)
ARR(X2_CX+XR, R3M, LD_X, R3M)           # X2 → Live Dashboard

# Live Dashboard
sp = RND(LD_X, R3T, LD_W, BH); T(sp, "Live\nDashboard", Pt(6.5), WHITE)
ARR(LD_X+LD_W, R3M, END_X, R3M)          # → End

# End oval
sp = OVL(END_X, R3T, END_W, BH, fill=NAVY); T(sp, "End", Pt(8.5), WHITE, bold=True)

# ── "Start new Chat" secondary oval (bottom-left of lane 3, like KPMG) ────
# In KPMG there's a small "Start new Chat" oval at the very bottom-left of lane 3
# that feeds back up into the Left-X. Mirror it exactly.
SNC_Y = R3T + BH + 0.32
sp = OVL(CX+0.10, SNC_Y, 0.95, BH-0.06, fill=NAVY)
T(sp, "Start new\nChat", Pt(6), WHITE)
ARR(CX+0.10+0.475, SNC_Y, XL1X, R3M+XR)   # up into Left-X bottom

# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
OUT = r"c:\Users\didar\Downloads\files (12)\telecom-complaint-system\TeleResolve_ToBe_Final.pptx"
prs.save(OUT)
print("Saved: " + OUT)
