# -*- coding: utf-8 -*-
"""
TeleResolve — To-Be Business Process (exact KPMG layout)
3 swim lanes + X-circle junctions + 3-row system lane
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ── Colours ────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1F, 0x39, 0x64)
BLUE  = RGBColor(0x41, 0x72, 0xC4)
CYAN  = RGBColor(0x00, 0xB0, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY  = RGBColor(0xBF, 0xBF, 0xBF)
DGREY = RGBColor(0x40, 0x40, 0x40)
LBBL  = RGBColor(0xEE, 0xF3, 0xFB)   # light blue lane BG
WBBL  = RGBColor(0xF9, 0xFB, 0xFF)   # near-white lane BG
RED   = RGBColor(0xC0, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Low-level shape helpers ────────────────────────────────────────────────
def _shape(mso, x, y, w, h, fill, lc=GREY, lw=Pt(0.75)):
    sp = slide.shapes.add_shape(mso, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if lc is None: sp.line.fill.background()
    else: sp.line.color.rgb = lc; sp.line.width = lw
    sp.text_frame.word_wrap = True
    return sp

def box(x, y, w, h, fill=BLUE, lc=WHITE, lw=Pt(0.75)):
    return _shape(1, x, y, w, h, fill, lc, lw)          # rectangle

def rbox(x, y, w, h, fill=BLUE, lc=WHITE, lw=Pt(0.75)):
    sp = _shape(5, x, y, w, h, fill, lc, lw)            # rounded rect
    sp.adjustments[0] = 0.09; return sp

def dmnd(x, y, w, h, fill=BLUE, lc=WHITE):
    return _shape(4, x, y, w, h, fill, lc, Pt(0.75))    # diamond

def ovl(x, y, w, h, fill=NAVY, lc=WHITE):
    return _shape(9, x, y, w, h, fill, lc, Pt(0.75))    # oval

def xcir(cx, cy, r=0.175):
    """Small dark circle with X — junction connector."""
    sp = ovl(cx-r, cy-r, r*2, r*2, fill=DGREY, lc=WHITE)
    wrt(sp, "X", Pt(6.5), WHITE, bold=True); return sp

def wrt(sp, text, sz=Pt(7), col=WHITE, bold=False, align=PP_ALIGN.CENTER):
    tf = sp.text_frame; p = tf.paragraphs[0]; p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text; r.font.size = sz; r.font.bold = bold; r.font.color.rgb = col

def lbl(x, y, w, h, text, sz=Pt(7), col=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align; r = p.add_run()
    r.text = text; r.font.size = sz; r.font.bold = bold; r.font.color.rgb = col

def arr(x1, y1, x2, y2, col=DGREY, w=Pt(1.0)):
    """Straight arrow connector."""
    con = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    con.line.color.rgb = col; con.line.width = w
    ln = con._element.spPr.find(qn('a:ln'))
    if ln is None: ln = etree.SubElement(con._element.spPr, qn('a:ln'))
    he = ln.find(qn('a:headEnd'))
    if he is None: he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type', 'arrow'); he.set('w', 'med'); he.set('len', 'med')
    te = ln.find(qn('a:tailEnd'))
    if te is None: te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'none')

# ═══════════════════════════════════════════════════════════════════════════
# LAYOUT CONSTANTS
# ═══════════════════════════════════════════════════════════════════════════
LLW  = 1.1     # lane label width
LX   = 0.15    # lane left edge
CX   = LX+LLW  # content area start x  (= 1.25)
CW   = 13.33 - LX - LLW - 0.18   # content width  (= 11.9)

# Swim lane y-ranges
L1_Y, L1_H = 0.62, 1.52    # Customer
L2_Y, L2_H = 2.14, 1.20    # Admin Approval
L3_Y, L3_H = 3.34, 3.68    # TeleResolve System  (to y=7.02)

BH = 0.44   # standard box height

# Vertical centres for each row
L1_CY = L1_Y + L1_H/2 - BH/2       # lane 1 box top  ≈ 0.84
L2_CY = L2_Y + L2_H/2 - BH/2       # lane 2 box top  ≈ 2.52
R1_CY = L3_Y + 0.55                 # row 1 box top   ≈ 3.89
R2_CY = L3_Y + 1.62                 # row 2 box top   ≈ 4.96
R3_CY = L3_Y + 2.72                 # row 3 box top   ≈ 6.06

# Row-1 box dimensions (5 boxes spread across content area)
R1_GAP = 0.13
R1_BW  = (CW - 5*R1_GAP) / 5       # ≈ 2.19  per box

# ═══════════════════════════════════════════════════════════════════════════
# BACKGROUND & SLIDE FRAME
# ═══════════════════════════════════════════════════════════════════════════
box(0, 0, 13.33, 7.5, fill=WHITE, lc=None)

# Header
box(0, 0, 13.33, 0.52, fill=NAVY, lc=None)
lbl(0.18, 0.06, 5, 0.22, "6. Business Architecture", sz=Pt(13), col=WHITE, bold=True)
lbl(0.18, 0.29, 9, 0.18, "Business Process (To-Be)  —  TeleResolve AI-Powered System",
    sz=Pt(8.5), col=CYAN)

# Footer
box(0, 7.12, 13.33, 0.38, fill=NAVY, lc=None)
lbl(0.2, 7.15, 12, 0.22,
    "© 2026 KPMG Assurance and Consulting Services LLP — an Indian Limited Liability Partnership and a member firm of the KPMG global organization of independent member firms affiliated with KPMG International Limited, a private English company limited by guarantee. All rights reserved.",
    sz=Pt(6.5), col=GREY)

# Diagram outer border
box(LX, 0.6, 13.33-LX*2, 6.52, fill=WHITE, lc=GREY, lw=Pt(1.0))

# ═══════════════════════════════════════════════════════════════════════════
# SWIM LANE BANDS
# ═══════════════════════════════════════════════════════════════════════════
for i, (name, ly, lh, bg) in enumerate([
        ("Customer",             L1_Y, L1_H, LBBL),
        ("Admin\nApproval",      L2_Y, L2_H, WBBL),
        ("TeleResolve\nSystem",  L3_Y, L3_H, LBBL),
]):
    # label column
    sp = box(LX, ly, LLW, lh, fill=NAVY, lc=WHITE, lw=Pt(0.5))
    wrt(sp, name, Pt(8.5), WHITE, bold=True)
    # content BG
    box(CX, ly, CW, lh, fill=bg, lc=GREY, lw=Pt(0.5))

# ═══════════════════════════════════════════════════════════════════════════
# LANE 1 — CUSTOMER
# ═══════════════════════════════════════════════════════════════════════════
# Start oval
cy1 = L1_CY                   # box top-y
sp = ovl(CX+0.05, cy1, 0.72, BH);  wrt(sp, "Start", Pt(8), WHITE, bold=True)

# arrow → diamond
arr(CX+0.77, cy1+BH/2, CX+1.1, cy1+BH/2)

# Diamond: Existing User?
DX = CX + 1.1                 # diamond left x
DW, DH = 1.05, 0.68
sp = dmnd(DX, cy1-0.12, DW, DH);  wrt(sp, "Existing\nUser?", Pt(6.5), WHITE)
D_CX = DX + DW/2              # diamond center x
D_CY = cy1 - 0.12 + DH/2     # diamond center y

# YES arc — from diamond-top up and right to Login
lbl(D_CX-0.08, cy1-0.18, 0.3, 0.18, "Yes", sz=Pt(6.5), col=DGREY)
arc_y = L1_Y + 0.06           # height of arc (near top of lane 1)
arr(D_CX, cy1-0.12, D_CX, arc_y)              # up from diamond top
arr(D_CX, arc_y, 12.28, arc_y)                # right along top
arr(12.28, arc_y, 12.28, cy1)                 # down to Login

# NO path — right from diamond to Register
lbl(DX+DW+0.04, D_CY-0.12, 0.25, 0.18, "No", sz=Pt(6.5), col=DGREY)
arr(DX+DW, D_CY, DX+DW+0.22, D_CY)

# Register
RGX = DX + DW + 0.22
sp = rbox(RGX, cy1, 1.05, BH);  wrt(sp, "Register", Pt(7), WHITE)
arr(RGX+1.05, cy1+BH/2, RGX+1.25, cy1+BH/2)

# Create Password
CPX = RGX + 1.25
sp = rbox(CPX, cy1, 1.15, BH);  wrt(sp, "Create\nPassword", Pt(6.5), WHITE)
arr(CPX+1.15, cy1+BH/2, CPX+1.35, cy1+BH/2)

# Generate OTP
OX = CPX + 1.35
sp = rbox(OX, cy1, 1.15, BH);  wrt(sp, "Generate\nOTP", Pt(6.5), WHITE)
arr(OX+1.15, cy1+BH/2, OX+1.35, cy1+BH/2)

# OTP → right → Login (merges with Yes arc at Login)
# Draw line from OTP right edge across to Login
LOGIN_X = 11.75
arr(OX+1.15, cy1+BH/2, LOGIN_X, cy1+BH/2)

# Login box (far right, both Yes and No paths merge here)
sp = rbox(LOGIN_X, cy1, 1.0, BH, fill=NAVY, lc=WHITE)
wrt(sp, "Login", Pt(8.5), WHITE, bold=True)

# ═══════════════════════════════════════════════════════════════════════════
# LANE 2 — ADMIN APPROVAL
# ═══════════════════════════════════════════════════════════════════════════
cy2 = L2_CY

# Arrow from Register bottom → down → Admin diamond
arr(RGX+0.52, cy1+BH, RGX+0.52, cy2+0.12)

# Admin Approval diamond (under Register)
ADX = RGX - 0.02
ADW, ADH = 1.1, 0.72
sp = dmnd(ADX, cy2, ADW, ADH);  wrt(sp, "Admin\nApproval", Pt(6.5), WHITE)
AD_CX = ADX + ADW/2
AD_CY = cy2 + ADH/2

# NO → left → Reject box
lbl(ADX-0.38, AD_CY-0.11, 0.25, 0.18, "No", sz=Pt(6.5), col=DGREY)
arr(ADX, AD_CY, CX+0.1, AD_CY)
sp = rbox(CX+0.1, cy2+0.14, 0.9, BH, fill=RED, lc=WHITE)
wrt(sp, "Reject &\nNotify User", Pt(6.5), WHITE)

# YES → right → up → merges into OTP flow (connects to space between OTP and Login)
lbl(ADX+ADW+0.04, AD_CY-0.11, 0.25, 0.18, "Yes", sz=Pt(6.5), col=DGREY)
YES_MERGE_X = OX + 0.57       # midpoint of OTP box, to merge into the flow
arr(ADX+ADW, AD_CY, YES_MERGE_X, AD_CY)      # right
arr(YES_MERGE_X, AD_CY, YES_MERGE_X, cy1+BH) # up to lane 1 bottom

# ═══════════════════════════════════════════════════════════════════════════
# LANE 3 — TELERESOLVE SYSTEM
# 3 rows of boxes with X-circle junctions (exact KPMG layout)
# ═══════════════════════════════════════════════════════════════════════════

# ── Row 1: 5 feature boxes (with left X and right X) ─────────────────────
R1Y  = R1_CY                  # box top y for row 1
R1MY = R1Y + BH/2             # box mid y for row 1

# Left X junction (entry from Login above)
XL1_CX = CX + 0.22
xcir(XL1_CX, R1MY)

# Arrow from Login bottom → down → left X
arr(LOGIN_X+0.5, cy1+BH, LOGIN_X+0.5, R1_Y := R1MY)

# Actually connect Login bottom to XL1 via a path: Login bottom → down to lane 3 → left to XL1
# Login center-x = LOGIN_X + 0.5
# XL1 center-x = XL1_CX
# We go: Login bottom (LOGIN_X+0.5, cy1+BH) → down to R1MY → left to XL1_CX (but XL1_CX < LOGIN_X)
# Better path: Login right → far right → down → left → XL1
# Per KPMG original: Login is far right, so the arrow goes straight down from Login into lane 3 near the right side
# Then inside lane 3 it connects. Let's mirror that:
# Arrow from Login bottom → straight down into lane 3 → left X is at left side
# But that's a long leftward line. Let me check the original:
# In KPMG original, the right side of lane 3 also has the X circle (right X of row 1),
# and Login seems to connect DOWN to the RIGHT X of row 1 (not the left X).
# Left X in row 1 also connects DOWN to left X in row 3.
# So Login → down → RIGHT X of row 1 → then row 1 boxes flow LEFT to RIGHT into right X... that's backwards.

# Actually re-reading: the flow in row 1 is LEFT→RIGHT. The LEFT X is the entry, the RIGHT X is the exit.
# Login connects down to LEFT X (which is on the left side of the system lane).
# This means Login (at x≈12.25) is far right in lane 1, and left X of system lane is at x≈1.45.
# They connect with a long leftward arrow or an L-shaped path.
# Looking at original KPMG: there seems to be a path from Login → down → far right in system lane → then connects to right side or...

# Let me reconsider the original layout. In the KPMG image:
# - Login is at the FAR RIGHT of lane 1 (x ≈ 12.5 inches)
# - The X circle on the LEFT of row 1 in the system lane is at x ≈ 1.5 inches
# - There appears to be a long arrow from Login going left somehow, OR
# - The Login arrow goes DOWN into the system lane from the right, and the right X circle is what connects from Login

# I think the correct interpretation is:
# - Login → arrow goes down into the system lane
# - This connects to the RIGHT X circle of row 1 (not the left)
# - The 5 boxes flow from LEFT to RIGHT, ending at RIGHT X
# - LEFT X of row 1 → down to LEFT X of row 3 → Start Chat → etc.
# This way, Login connects at the RIGHT side, flow goes left→right,
# and the left X is where row 3 starts.

# But that would mean row 1 flows right→left (from right X to left X going LEFT)
# which is unusual for flowcharts.

# Let me re-examine: In the KPMG Control Iris Tool diagram:
# Row 1 boxes: Create New Project | Open Existing Project | Recent Projects | Open Repository | Upload Documents
# These appear to be MENU OPTIONS (not sequential), so the arrows between them might show navigation options
# The X circle on left = entry junction (split: either enter row 1 features OR go to row 3 chat)
# The X circle on right = exit junction of row 1 features

# So the flow is:
# Login → (somewhere) → enters system → Left X (row1/row3 junction)
# At Left X:
#   Option A: go RIGHT through 5 feature boxes → Right X → down → row 2 (Select Phase/View Dashboard)
#   Option B: go DOWN to Left X of row 3 → Start Chat → X → Run Audit → Live Dashboard → End

# This makes sense! Login → left X junction (split point):
# Path A: menu features → workflow
# Path B: direct to chat/audit

# For TeleResolve:
# Login → Left X junction:
# Path A: (row 1) Submit Complaint | View Tickets | Recent Sessions | Network AI | Upload Evidence → Right X → down → Select Issue Category → View Dashboard
# Path B: (row 3) Start AI Chat → X → Run Diagnosis → Live Dashboard → End

# So I need to connect Login to the LEFT X circle somehow
# In KPMG original: there's probably an arrow from Login going all the way left to the left X
# Or Login connects to a point in lane 3 that then connects to the left X

# Let me just draw it naturally: Login bottom → long arrow going left down into the system lane to the left X
# This is what the KPMG diagram does with a connecting line that goes left

# Actually I'll look at this differently. The "Yes" label at the VERY TOP of the KPMG slide
# seems to indicate a top-level arc. The Login box is indeed at the right side.
# After Login, an arrow goes DOWN from Login box into the system lane near the left side.
# This would require the arrow to go: Login_bottom → down a bit → left (horizontal) → into left X of row1

# Let me draw: Login → down → then a line going left at the top of lane 3 → into left X

# I'll draw it as: Login bottom → down to lane 3 → horizontal left line → Left X circle
# This is the natural KPMG-style routing

# Coordinates:
# Login is at x=LOGIN_X (11.75) to x=LOGIN_X+1.0 (12.75), so center-x = 12.25
# Left X of row 1 is at CX+0.22 = 1.47
# Top of lane 3 = L3_Y = 3.34

# Draw the path:
# 1. Login bottom-center (12.25, cy1+BH) → down to (12.25, L3_Y+0.3)
# 2. (12.25, L3_Y+0.3) → left to (XL1_CX, L3_Y+0.3)   [horizontal line at top of lane 3]
# 3. (XL1_CX, L3_Y+0.3) → down to (XL1_CX, R1MY)      [connect to left X circle]
# But python-pptx connectors are only straight lines between two points, not bent.
# So I need two separate straight arrows to create the L-shape.

arr(LOGIN_X+0.5, cy1+BH, LOGIN_X+0.5, L3_Y+0.26)   # Login → down
arr(LOGIN_X+0.5, L3_Y+0.26, XL1_CX, L3_Y+0.26)      # right to left long line
arr(XL1_CX, L3_Y+0.26, XL1_CX, R1MY-0.175)          # down into X circle

# Row 1: 5 feature boxes
r1_labels = [
    "Submit\nComplaint",
    "View Active\nTickets",
    "Recent\nSessions",
    "Network AI\nChat",
    "Upload\nEvidence",
]
r1_starts = []
rx = XL1_CX + 0.22 + 0.175   # start after left X
arr(XL1_CX+0.175, R1MY, rx, R1MY)  # X → first box

for i, lbl_txt in enumerate(r1_labels):
    sp = rbox(rx, R1Y, R1_BW, BH, fill=BLUE, lc=WHITE)
    wrt(sp, lbl_txt, Pt(6.5), WHITE)
    r1_starts.append(rx)
    if i < len(r1_labels)-1:
        arr(rx+R1_BW, R1MY, rx+R1_BW+R1_GAP, R1MY)
        rx += R1_BW + R1_GAP
    else:
        rx += R1_BW

# Right X circle (after last box in row 1)
XR1_CX = rx + 0.22
arr(rx, R1MY, XR1_CX-0.175, R1MY)
xcir(XR1_CX, R1MY)

# ── Row 2: Select Issue Category → View Dashboard ─────────────────────────
R2Y  = R2_CY
R2MY = R2Y + BH/2

# Right X → down → row 2
arr(XR1_CX, R1MY+0.175, XR1_CX, R2MY)
arr(XR1_CX, R2MY, XR1_CX-1.55, R2MY)   # left to Select box

SIC_X = XR1_CX - 1.55 - 1.3
arr(XR1_CX-1.55, R2MY, SIC_X+1.3, R2MY)
sp = rbox(SIC_X, R2Y, 1.3, BH, fill=BLUE, lc=WHITE)
wrt(sp, "Select Issue\nCategory", Pt(6.5), WHITE)
arr(SIC_X, R2MY, SIC_X-0.15, R2MY)

VD_X = SIC_X - 0.15 - 1.3
arr(SIC_X-0.15, R2MY, VD_X+1.3, R2MY)
sp = rbox(VD_X, R2Y, 1.3, BH, fill=BLUE, lc=WHITE)
wrt(sp, "View\nDashboard", Pt(6.5), WHITE)

# ── Row 3: Start AI Chat → X → Run Network Diagnosis → Live Dashboard → End
R3Y  = R3_CY
R3MY = R3Y + BH/2

# Left X circle row 3 (connects DOWN from left X row 1)
XL3_CX = XL1_CX
xcir(XL3_CX, R3MY)
arr(XL1_CX, R1MY+0.175, XL3_CX, R3MY-0.175)   # left X row1 → down → left X row3

arr(XL3_CX+0.175, R3MY, XL3_CX+0.4, R3MY)

# Start AI Chat (oval, like KPMG "Start new Chat" oval)
SCH_X = XL3_CX + 0.4
sp = ovl(SCH_X, R3Y, 1.25, BH, fill=NAVY, lc=WHITE)
wrt(sp, "Start AI\nChat Session", Pt(6.5), WHITE)
arr(SCH_X+1.25, R3MY, SCH_X+1.45, R3MY)

# X circle between Start and Diagnose
XM1_CX = SCH_X + 1.45 + 0.175
xcir(XM1_CX, R3MY)
arr(SCH_X+1.45, R3MY, XM1_CX-0.175, R3MY)
arr(XM1_CX+0.175, R3MY, XM1_CX+0.35, R3MY)

# Run Network Diagnosis box
RND_X = XM1_CX + 0.35
sp = rbox(RND_X, R3Y, 1.4, BH, fill=BLUE, lc=WHITE)
wrt(sp, "Run Network\nDiagnosis", Pt(6.5), WHITE)
arr(RND_X+1.4, R3MY, RND_X+1.6, R3MY)

# X circle between Diagnose and Dashboard
XM2_CX = RND_X + 1.6 + 0.175
xcir(XM2_CX, R3MY)
arr(RND_X+1.6, R3MY, XM2_CX-0.175, R3MY)
arr(XM2_CX+0.175, R3MY, XM2_CX+0.35, R3MY)

# Live Dashboard box
LDB_X = XM2_CX + 0.35
sp = rbox(LDB_X, R3Y, 1.35, BH, fill=BLUE, lc=WHITE)
wrt(sp, "Live\nDashboard", Pt(6.5), WHITE)
arr(LDB_X+1.35, R3MY, LDB_X+1.55, R3MY)

# End oval
END_X = LDB_X + 1.55
sp = ovl(END_X, R3Y, 0.78, BH, fill=NAVY, lc=WHITE)
wrt(sp, "End", Pt(8), WHITE, bold=True)

# ── View Dashboard → down → connects to Live Dashboard (or row 3 X) ───────
# View Dashboard right edge connects down and right into row 3 area
arr(VD_X, R2MY, VD_X-0.3, R2MY)
arr(VD_X-0.3, R2MY, VD_X-0.3, XM2_CX)   # reuse XM2_CX as y target (same level)
arr(VD_X-0.3, R3MY, XM2_CX-0.175, R3MY)  # connect to second X circle in row 3

# ── "Start new Chat" from left of lane 3 (small oval, like KPMG) ──────────
# Mirror the KPMG "Start new Chat" box which is on the far LEFT of row 3 BELOW the left X
SNC_X = CX + 0.12
SNC_Y = R3Y + BH + 0.35
sp = ovl(SNC_X, SNC_Y, 0.9, BH-0.06, fill=NAVY, lc=WHITE)
wrt(sp, "Start\nnew Chat", Pt(6), WHITE)
arr(SNC_X+0.45, SNC_Y, XL3_CX, R3MY+0.175)   # connects up to left X row 3

# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
OUT = r"c:\Users\didar\Downloads\files (12)\telecom-complaint-system\TeleResolve_ToBe.pptx"
prs.save(OUT)
print("Saved: " + OUT)
